import io
import re
from typing import Dict, List, Optional, Tuple, Union

import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# --------------- Streamlit page config (must be first Streamlit call) ---------------
st.set_page_config(
    page_title="Asset Allocation - Tabela Única + PPTX (Layout Estruturado)",
    layout="centered",
)

# --------------- Simple Authentication ---------------
if not st.session_state.get("authenticated", False):
    st.warning("Please enter the password on the Home page first.")
    st.write("Status: Não Autenticado")
    st.stop()

st.write("Autenticado")

# --------------- Constants (Ceres Wealth style guide) ---------------

FONT_NAME = "Century Gothic"

# Colors (Ceres Wealth palette)
TITLE_COLOR = "#825120"  # header/subheader brown
DEFAULT_ZEBRA_HEX = "#EFEFEF"
HEADER_FG_HEX = "#FFFFFF"

# Derived styling
SECTION_LIGHT_ALPHA = 0.14  # light tint for subsection rows on web/PPT
COLUMN_GAP_PX = 14  # web visual gap between groups (header only)

# RGB tuples for PPTX
HEADER_BG_RGB = (0x82, 0x51, 0x20)  # #825120

# A4 Size (vertical)
SLIDE_WIDTH = Inches(8.27)
SLIDE_HEIGHT = Inches(11.69)

# Defaults for sizes and paddings (inches and points)
HDR_ROW_H_IN = 0.28
SUBHDR_ROW_H_IN = 0.26
BODY_ROW_H_IN = 0.24
MARGIN_IN = 0.5

# Canonical structure (order + indent level) for Page 1
STRUCTURE: List[Tuple[str, int]] = [
    ("Brasil", 0),
    ("Renda Fixa", 1),
    ("Pós Fixado", 2),
    ("Pré-Fixado", 2),
    ("Inflação", 2),
    ("Multimercado", 0),
    ("No Brasil", 1),
    ("No Exterior", 1),
    ("Renda Variável", 0),
    ("Ações", 1),
    ("Fundo de Ações", 1),
    ("FIIs | FIAgro | FIInfra", 1),
    ("Alternativos", 0),
    ("Global", 0),
    ("Renda Fixa", 1),
    ("Renda Variável", 1),
    ("TOTAL", 0),
]

# Labels and groupings (used by both pages)
TOP_LEVEL_LABELS = {
    "Brasil",
    "Multimercado",
    "Multimercados",
    "Renda Variável",
    "Alternativos",
    "Global",
    "TOTAL",
    "Fundos de Renda Fixa",
    "Renda Fixa - Títulos Públicos",
    "Renda Fixa - Crédito Bancário",
    "Renda Fixa - Crédito Privado",
}

TOP_SECTION_LABELS = {"Brasil", "Global", "TOTAL"}

# Subsection rows that get the light brown tint and bold
SUBSECTION_LABELS = {
    "Renda Fixa",
    "Multimercado",
    "Renda Variável",
    "Alternativos",
    "Fundos de Renda Fixa",
    "Multimercados",
    "Renda Fixa - Títulos Públicos",
    "Renda Fixa - Crédito Bancário",
    "Renda Fixa - Crédito Privado",
}

# Column gap boundaries (after these column indices, add a web-only header gap)
# 0: Asset Allocation
# 1: Cons Alvo, 2: Cons Faixa, 3: Mod Alvo, 4: Mod Faixa, 5: Sof Alvo, 6: Sof Faixa
GAP_AFTER_COLS = {0, 2, 4}

# --------------- Utilities ---------------


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.strip().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def hex_to_rgb_tuple(hex_str: str) -> Tuple[int, int, int]:
    h = hex_str.strip().lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def hex_to_rgba_str(hex_str: str, alpha: float) -> str:
    h = hex_str.strip().lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    a = min(max(alpha, 0.0), 1.0)
    return f"rgba({r}, {g}, {b}, {a})"


def normalize_text(text: str) -> str:
    return (
        text.replace("\u00A0", " ")
        .replace("\r\n", "\n")
        .replace("\r", "\n")
        .strip()
    )


def sanitize_line(line: str) -> str:
    return re.sub(r"\s+", " ", line.strip())


def split_cells(line: str) -> List[str]:
    parts = re.split(r"\s*\|\s*|\t+| {2,}", line)
    return [p.strip() for p in parts if p and p.strip()]


def is_percent(token: str) -> bool:
    return bool(re.match(r"^\d{1,3}(?:[.,]\d+)?%$", token))


def is_range(token: str) -> bool:
    return bool(re.match(r"^\d{1,3}%\s*-\s*\d{1,3}%$", token))


def parse_label_and_values(line: str) -> Optional[Tuple[str, List[str]]]:
    tokens = split_cells(line)
    if not tokens:
        return None
    idx = 0
    while idx < len(tokens) and not (is_percent(tokens[idx]) or is_range(tokens[idx])):
        idx += 1
    label_tokens = tokens[:idx]
    value_tokens = tokens[idx:]
    if not label_tokens and not value_tokens:
        return None
    label_joiner = " | " if "|" in line else " "
    label = label_joiner.join(label_tokens).strip()
    return (label, value_tokens)


# --------------- Parsing for the single table ---------------


def extract_profile_pairs(vals: List[str]) -> Tuple[str, str, str, str, str, str]:
    vals = [v for v in vals if v]
    i = 0
    out: List[Tuple[str, str]] = []
    for _ in range(3):
        alvo, faixa = "", ""
        if i < len(vals) and is_percent(vals[i]):
            alvo = vals[i]
            i += 1
        if i < len(vals) and is_range(vals[i]):
            faixa = vals[i]
            i += 1
        out.append((alvo, faixa))
    while len(out) < 3:
        out.append(("", ""))
    c_a, c_f = out[0]
    m_a, m_f = out[1]
    s_a, s_f = out[2]
    return c_a, c_f, m_a, m_f, s_a, s_f


def parse_asset_allocation_table(lines: List[str]) -> pd.DataFrame:
    rows = []
    for raw in lines:
        s_clean = sanitize_line(raw)
        s_lower = s_clean.lower()
        if not s_clean:
            continue
        if "asset allocation" in s_lower:
            continue
        if "alvo" in s_lower and "faixa" in s_lower:
            continue

        parsed = parse_label_and_values(raw)
        if not parsed:
            continue

        label, vals = parsed
        if not label and not vals:
            continue

        c_a, c_f, m_a, m_f, s_a, s_f = extract_profile_pairs(vals)
        rows.append(
            {
                "Item": label,
                "Cons Alvo": c_a,
                "Cons Faixa": c_f,
                "Mod Alvo": m_a,
                "Mod Faixa": m_f,
                "Sof Alvo": s_a,
                "Sof Faixa": s_f,
            }
        )

    cols = [
        "Item",
        "Cons Alvo",
        "Cons Faixa",
        "Mod Alvo",
        "Mod Faixa",
        "Sof Alvo",
        "Sof Faixa",
    ]
    if not rows:
        return pd.DataFrame(columns=cols)
    return pd.DataFrame(rows, columns=cols)


def reorder_to_structure(
    df_simple: pd.DataFrame,
) -> Tuple[pd.DataFrame, List[int]]:
    label_to_indices: Dict[str, List[int]] = {}
    for i, label in enumerate(df_simple["Item"].tolist()):
        label_to_indices.setdefault(label, []).append(i)

    ordered_rows: List[Dict[str, str]] = []
    indents: List[int] = []

    for label, indent in STRUCTURE:
        idx_list = label_to_indices.get(label, [])
        if idx_list:
            use_idx = idx_list.pop(0)
            row = df_simple.iloc[use_idx].to_dict()
        else:
            row = {
                "Item": label,
                "Cons Alvo": "",
                "Cons Faixa": "",
                "Mod Alvo": "",
                "Mod Faixa": "",
                "Sof Alvo": "",
                "Sof Faixa": "",
            }
        ordered_rows.append(row)
        indents.append(indent)

    ordered_df = pd.DataFrame(ordered_rows, columns=df_simple.columns)
    return ordered_df, indents


def to_multiindex_view(df_simple: pd.DataFrame) -> pd.DataFrame:
    cols = pd.MultiIndex.from_tuples(
        [
            ("Asset Allocation", ""),
            ("Conservador", "Alvo (%)"),
            ("Conservador", "Faixa Tolerada (%)"),
            ("Moderado", "Alvo (%)"),
            ("Moderado", "Faixa Tolerada (%)"),
            ("Sofisticado", "Alvo (%)"),
            ("Sofisticado", "Faixa Tolerada (%)"),
        ]
    )
    df_mx = pd.DataFrame(
        {
            ("Asset Allocation", ""): df_simple["Item"].tolist(),
            ("Conservador", "Alvo (%)"): df_simple["Cons Alvo"].tolist(),
            ("Conservador", "Faixa Tolerada (%)"): df_simple["Cons Faixa"].tolist(),
            ("Moderado", "Alvo (%)"): df_simple["Mod Alvo"].tolist(),
            ("Moderado", "Faixa Tolerada (%)"): df_simple["Mod Faixa"].tolist(),
            ("Sofisticado", "Alvo (%)"): df_simple["Sof Alvo"].tolist(),
            ("Sofisticado", "Faixa Tolerada (%)"): df_simple["Sof Faixa"].tolist(),
        }
    )
    df_mx = df_mx[cols]
    return df_mx


def parse_single_table_text(raw_text: str) -> Tuple[pd.DataFrame, List[int]]:
    text = normalize_text(raw_text)
    lines = text.split("\n")
    df_simple = parse_asset_allocation_table(lines)
    ordered_df, indents = reorder_to_structure(df_simple)
    return ordered_df, indents


def parse_free_table_text(raw_text: str) -> Tuple[pd.DataFrame, List[int]]:
    text = normalize_text(raw_text)
    lines = text.split("\n")
    df_simple = parse_asset_allocation_table(lines)
    indents = [0] * len(df_simple)
    return df_simple, indents


# --------------- Markdown helpers (optional preview/export) ---------------


def escape_pipes(s: str) -> str:
    return s.replace("|", r"\|" )


def flatten_columns_labels_for_text(cols) -> List[str]:
    labels: List[str] = []
    if isinstance(cols, pd.MultiIndex):
        for a, b in cols:
            if a and b:
                labels.append(f"{a} / {b}")
            else:
                labels.append(a or b or "")
    else:
        labels = [str(c) for c in cols]
    return labels


def df_to_markdown(df: pd.DataFrame) -> str:
    col_labels = flatten_columns_labels_for_text(df.columns)
    header = "| " + " | ".join(escape_pipes(c) for c in col_labels) + " |"
    sep = "| " + " | ".join("---" for _ in col_labels) + " |"
    lines = [header, sep]
    for _, r in df.iterrows():
        row = ["" if (isinstance(v, float) and pd.isna(v)) else str(v) for v in r]
        lines.append("| " + " | ".join(escape_pipes(x) for x in row) + " |")
    return "\n".join(lines)


# --------------- Styling helpers ---------------


def is_subsection_label(label: str) -> bool:
    lab = label.strip().lower()
    if lab in {s.lower() for s in SUBSECTION_LABELS}:
        return True
    return lab.startswith("renda fixa - ")


# --------------- PPTX helpers ---------------


def set_paragraph(
    paragraph,
    text: str,
    size: int,
    bold: bool = False,
    color: Optional[str] = None,
    align: Optional[int] = None,
):
    if hasattr(paragraph, "clear"):
        paragraph.clear()
    else:
        paragraph.text = ""
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = FONT_NAME
    font.size = Pt(size)
    font.bold = bold
    if color:
        font.color.rgb = hex_to_rgb(color)
    if align is not None:
        paragraph.alignment = align


def set_cell_fill_with_alpha(cell, rgb_tuple: Tuple[int, int, int], alpha: float):
    from pptx.oxml.xmlchemy import OxmlElement

    try:
        tcPr = cell._tc.get_or_add_tcPr()
        for fill in tcPr.findall(qn("a:solidFill")):
            tcPr.remove(fill)
        solidFill = OxmlElement("a:solidFill")
        srgbClr = OxmlElement("a:srgbClr")
        hex_val = "{:02X}{:02X}{:02X}".format(
            rgb_tuple[0], rgb_tuple[1], rgb_tuple[2]
        )
        srgbClr.set("val", hex_val)
        alpha_elem = OxmlElement("a:alpha")
        alpha_elem.set("val", str(int(alpha * 100000)))
        srgbClr.append(alpha_elem)
        solidFill.append(srgbClr)
        tcPr.append(solidFill)
    except Exception:
        pass


def set_cell_text(
    cell,
    text: str,
    size: int,
    bold: bool = False,
    font_color_hex: Optional[str] = None,
    align: int = PP_ALIGN.CENTER,
):
    cell.text_frame.clear()
    cell.text_frame.word_wrap = True
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    color_hex = font_color_hex or "#000000"
    run.font.color.rgb = hex_to_rgb(color_hex)


def apply_borderless_theme_to_table(table) -> None:
    from pptx.oxml.xmlchemy import OxmlElement

    try:
        tbl = table._tbl
        for tc in tbl.iter(qn("a:tc")):
            tcPr = tc.get_or_add_tcPr()
            for tag in [qn("a:lnL"), qn("a:lnR"), qn("a:lnT"), qn("a:lnB")]:
                for elem in tcPr.findall(tag):
                    tcPr.remove(elem)
            for pos in ["lnL", "lnR", "lnT", "lnB"]:
                ln = OxmlElement(f"a:{pos}")
                noFill = OxmlElement("a:noFill")
                ln.append(noFill)
                tcPr.append(ln)
    except Exception:
        pass


def add_background(
    slide,
    prs: Presentation,
    bg_source: Optional[Union[str, io.BytesIO]] = None,
):
    if bg_source is None:
        return
    try:
        slide.shapes.add_picture(
            bg_source, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
    except Exception:
        pass


def add_cover_slide(
    prs: Presentation,
    title: str,
    bg_source: Optional[Union[str, io.BytesIO]] = None,
    font_color_hex: str = "#000000",
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs, bg_source)

    left = Inches(MARGIN_IN)
    top = Inches(3.2)
    width = prs.slide_width - Inches(2 * MARGIN_IN)
    height = Inches(1.8)
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    set_paragraph(p, title, 44, True, font_color_hex, PP_ALIGN.CENTER)


def add_df_table(
    slide,
    df: pd.DataFrame,
    left,
    top,
    width,
    hdr_size: int,
    body_size: int,
    font_color_hex: str,
    zebra_rgb: Tuple[int, int, int],
    zebra_alpha: float,
    indent_levels: Optional[List[int]] = None,
    first_col_pct: float = 30.0,
    subheader_extra: Optional[List[str]] = None,
    top_blacklist: Optional[List[str]] = None,
):
    data_rows = len(df)
    cols_n = len(df.columns)
    rows_n = 2 + data_rows

    hdr_h = Inches(HDR_ROW_H_IN)
    subhdr_h = Inches(SUBHDR_ROW_H_IN)
    row_h = Inches(BODY_ROW_H_IN)
    table_h = hdr_h + subhdr_h + row_h * max(1, data_rows)

    shape = slide.shapes.add_table(rows_n, cols_n, left, top, width, table_h)
    table = shape.table
    apply_borderless_theme_to_table(table)

    table.rows[0].height = hdr_h
    table.rows[1].height = subhdr_h
    for i in range(2, rows_n):
        table.rows[i].height = row_h

    first_fraction = max(0.1, min(first_col_pct / 100.0, 0.6))
    width_emu = int(width)
    first_w = int(width_emu * first_fraction)
    rest_total = width_emu - first_w
    per_rest = int(rest_total / max(1, cols_n - 1))
    for j in range(cols_n):
        table.columns[j].width = first_w if j == 0 else per_rest

    table.cell(0, 0).merge(table.cell(1, 0))
    table.cell(0, 1).merge(table.cell(0, 2))
    table.cell(0, 3).merge(table.cell(0, 4))
    table.cell(0, 5).merge(table.cell(0, 6))

    top_labels = ["Asset Allocation", "Conservador", "Moderado", "Sofisticado"]
    top_positions = [(0, 0), (0, 1), (0, 3), (0, 5)]
    for (r, c), label in zip(top_positions, top_labels):
        cell = table.cell(r, c)
        set_cell_fill_with_alpha(cell, HEADER_BG_RGB, 1.0)
        set_cell_text(cell, label, hdr_size, True, HEADER_FG_HEX, PP_ALIGN.CENTER)

    second_labels = [
        "Alvo (%)",
        "Faixa Tolerada (%)",
        "Alvo (%)",
        "Faixa Tolerada (%)",
        "Alvo (%)",
        "Faixa Tolerada (%)",
    ]
    for j, label in enumerate(second_labels, start=1):
        cell = table.cell(1, j)
        set_cell_fill_with_alpha(cell, HEADER_BG_RGB, 1.0)
        set_cell_text(cell, label, hdr_size, True, HEADER_FG_HEX, PP_ALIGN.CENTER)

    indents = indent_levels or [0] * data_rows
    extra_sub = set(subheader_extra or [])
    top_bl = set(top_blacklist or [])
    top_level_lower = {s.lower() for s in TOP_LEVEL_LABELS}
    top_section_lower = {s.lower() for s in TOP_SECTION_LABELS}
    extra_sub_lower = {s.lower() for s in extra_sub}
    top_bl_lower = {s.lower() for s in top_bl}

    for i in range(data_rows):
        is_even = (i + 1) % 2 == 0
        label = "" if pd.isna(df.iloc[i, 0]) else str(df.iloc[i, 0]).strip()
        label_lower = label.lower()
        row_idx = 2 + i
        is_top_default = label_lower in top_section_lower
        row_is_top_section = is_top_default and (label_lower not in top_bl_lower)
        row_is_subsection = is_subsection_label(label) or (label_lower in extra_sub_lower)

        for j in range(cols_n):
            val = "" if pd.isna(df.iloc[i, j]) else str(df.iloc[i, j])
            cell = table.cell(row_idx, j)

            if row_is_top_section:
                set_cell_fill_with_alpha(cell, HEADER_BG_RGB, 1.0)
                fg_hex = HEADER_FG_HEX
            elif row_is_subsection:
                set_cell_fill_with_alpha(cell, HEADER_BG_RGB, SECTION_LIGHT_ALPHA)
                fg_hex = font_color_hex
            else:
                if is_even:
                    set_cell_fill_with_alpha(cell, zebra_rgb, zebra_alpha)
                else:
                    cell.fill.background()
                fg_hex = font_color_hex

            bold_flag = row_is_top_section or row_is_subsection or (
                label_lower in top_level_lower
            )
            set_cell_text(
                cell,
                val,
                body_size,
                bold=bold_flag,
                font_color_hex=fg_hex,
                align=PP_ALIGN.CENTER,
            )

    return {"shape": shape, "height": table_h}


def add_single_table_slide(
    prs: Presentation,
    page_title: str,
    df: pd.DataFrame,
    indent_levels: List[int],
    title_font_size: int,
    hdr_size: int,
    body_size: int,
    title_top_offset: float,
    title_padding: float,
    bg_source: Optional[Union[str, io.BytesIO]] = None,
    font_color_hex: str = "#000000",
    zebra_color_hex: str = DEFAULT_ZEBRA_HEX,
    zebra_alpha: float = 1.0,
    first_col_pct: float = 30.0,
    subheader_extra: Optional[List[str]] = None,
    top_blacklist: Optional[List[str]] = None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs, bg_source)

    margin = Inches(MARGIN_IN)
    content_w = prs.slide_width - 2 * margin
    zebra_rgb = hex_to_rgb_tuple(zebra_color_hex)

    title_top = Inches(title_top_offset)
    tb_title = slide.shapes.add_textbox(margin, title_top, content_w, Inches(0.7))
    p = tb_title.text_frame.paragraphs[0]
    p.text = ""
    set_paragraph(p, page_title, title_font_size, True, font_color_hex, PP_ALIGN.CENTER)

    current_top = title_top + Inches(title_padding)

    add_df_table(
        slide=slide,
        df=df,
        left=margin,
        top=current_top,
        width=content_w,
        hdr_size=hdr_size,
        body_size=body_size,
        font_color_hex=font_color_hex,
        zebra_rgb=zebra_rgb,
        zebra_alpha=zebra_alpha,
        indent_levels=indent_levels,
        first_col_pct=first_col_pct,
        subheader_extra=subheader_extra,
        top_blacklist=top_blacklist,
    )


def build_pptx(
    nome_doc: str,
    pages: List[Dict[str, Union[str, pd.DataFrame, List[int], List[str]]]],
    title_font_size: int,
    hdr_size: int,
    body_size: int,
    title_top_offset: float,
    title_padding: float,
    bg_source: Optional[Union[str, io.BytesIO]] = None,
    font_color_hex: str = "#000000",
    zebra_color_hex: str = DEFAULT_ZEBRA_HEX,
    zebra_alpha: float = 1.0,
    first_col_pct: float = 30.0,
) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    if nome_doc.strip():
        add_cover_slide(prs, nome_doc.strip(), bg_source, font_color_hex)

    for page in pages:
        add_single_table_slide(
            prs=prs,
            page_title=str(page["nome"]),
            df=page["df"],
            indent_levels=page["indent"],  # type: ignore[index]
            title_font_size=title_font_size,
            hdr_size=hdr_size,
            body_size=body_size,
            title_top_offset=title_top_offset,
            title_padding=title_padding,
            bg_source=bg_source,
            font_color_hex=font_color_hex,
            zebra_color_hex=zebra_color_hex,
            zebra_alpha=zebra_alpha,
            first_col_pct=first_col_pct,
            subheader_extra=page.get("subheader_extra", []),  # type: ignore[arg-type]
            top_blacklist=page.get("top_blacklist", []),  # type: ignore[arg-type]
        )

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


# --------------- Streamlit helpers ---------------


def style_table(
    df: pd.DataFrame,
    zebra_hex: str,
    zebra_alpha: float,
    indent_levels: List[int],
    first_col_pct: float,
    subheader_extra: Optional[List[str]] = None,
    top_blacklist: Optional[List[str]] = None,
) -> pd.io.formats.style.Styler:
    zebra_bg_css = hex_to_rgba_str(zebra_hex, zebra_alpha)
    subsection_bg_css = hex_to_rgba_str(TITLE_COLOR, SECTION_LIGHT_ALPHA)

    n_cols = df.shape[1]
    first_col = df.columns[0]
    rest_pct = max(0.0, (100.0 - first_col_pct) / max(1, n_cols - 1))

    extra_sub = set(subheader_extra or [])
    top_bl = set(top_blacklist or [])
    top_level_lower = {s.lower() for s in TOP_LEVEL_LABELS}
    top_section_lower = {s.lower() for s in TOP_SECTION_LABELS}
    extra_sub_lower = {s.lower() for s in extra_sub}
    top_bl_lower = {s.lower() for s in top_bl}

    base_styles = [
        {
            "selector": "th",
            "props": [
                ("background-color", TITLE_COLOR),
                ("color", "white"),
                ("font-weight", "bold"),
                ("font-family", "'Century Gothic', sans-serif"),
                ("border", "none"),
                ("text-align", "center"),
            ],
        },
        {
            "selector": "td",
            "props": [
                ("font-family", "'Century Gothic', sans-serif"),
                ("border", "none"),
                ("text-align", "center"),
            ],
        },
        {"selector": "table", "props": [("border-collapse", "collapse"), ("border", "none")]},
    ]

    for k in sorted(GAP_AFTER_COLS):
        nth = k + 1
        base_styles.append(
            {
                "selector": f"th.col_heading:nth-child({nth})",
                "props": [("border-right", f"{COLUMN_GAP_PX}px solid #FFFFFF")],
            }
        )
        base_styles.append(
            {
                "selector": f"th.col_heading.level1:nth-child({nth})",
                "props": [("border-right", f"{COLUMN_GAP_PX}px solid #FFFFFF")],
            }
        )

    base_styles.append(
        {"selector": "th.col_heading:nth-child(1)", "props": [("width", f"{first_col_pct}%"), ("min-width", "120px")]}
    )
    base_styles.append(
        {"selector": "td:nth-child(1)", "props": [("width", f"{first_col_pct}%"), ("min-width", "120px")]}
    )
    for j in range(2, n_cols + 1):
        base_styles.append({"selector": f"th.col_heading:nth-child({j})", "props": [("width", f"{rest_pct}%")]})
        base_styles.append({"selector": f"td:nth-child({j})", "props": [("width", f"{rest_pct}%")]})

    styler = df.style.set_table_styles(base_styles)

    def row_background(row):
        label = str(row[first_col]).strip().lower()
        is_top_default = label in top_section_lower
        row_is_top_section = is_top_default and (label not in top_bl_lower)
        row_is_subsection = is_subsection_label(str(row[first_col])) or (label in extra_sub_lower)
        if row_is_top_section:
            return [f"background-color: {TITLE_COLOR}; color: white; font-weight: bold;"] * n_cols
        if row_is_subsection:
            return [f"background-color: {subsection_bg_css}; font-weight: bold;"] * n_cols
        color = zebra_bg_css if row.name % 2 == 0 else "#FFFFFF"
        return [f"background-color: {color}"] * n_cols

    def bold_top_rows(row):
        label = str(row[first_col]).strip().lower()
        is_group = (
            indent_levels[row.name] == 0
            or label == "total"
            or label in top_level_lower
            or is_subsection_label(str(row[first_col]))
            or (label in extra_sub_lower)
        )
        weight = "bold" if is_group else "normal"
        return [f"font-weight: {weight};"] * n_cols

    def column_gaps(row):
        out = []
        for j in range(n_cols):
            if j in GAP_AFTER_COLS:
                out.append(f"border-right: {COLUMN_GAP_PX}px solid #FFFFFF;")
            else:
                out.append("")
        return out

    styler = styler.apply(row_background, axis=1)
    styler = styler.apply(bold_top_rows, axis=1)
    styler = styler.apply(column_gaps, axis=1)
    return styler


# --------------- Streamlit UI ---------------

st.markdown(
    "<h2 style='text-align: center;'>Parser: Tabela Única de Asset Allocation "
    "(visual estruturado) + Export PPTX</h2>",
    unsafe_allow_html=True,
)

st.markdown("<br><br>", unsafe_allow_html=True)

st.markdown(
    "- Cole abaixo os textos das tabelas de Asset Allocation.\n"
    "- O parser aceita colunas separadas por pipes (|), TABs, ou 2+ espaços.\n"
    "- Faixas no formato '10% - 25%' são reconhecidas.\n"
    "- As linhas são estilizadas para seguir o layout e as cores da referência."
)

if "pages" not in st.session_state:
    st.session_state.pages = []

st.subheader("Configurações")

cfg_col1, cfg_col2, cfg_col3, cfg_col4, cfg_col5 = st.columns(5)

with cfg_col1:
    nome_doc = st.text_input("Nome do Documento (capa PPT)", value="")

# Defaults updated per your screenshot:
with cfg_col2:
    title_font_size = st.slider("Tamanho título PPT", min_value=10, max_value=36, value=13, step=1)

with cfg_col3:
    hdr_font_size = st.slider("Tamanho fonte headers", min_value=8, max_value=16, value=9, step=1)

with cfg_col4:
    body_font_size = st.slider("Tamanho fonte corpo", min_value=7, max_value=14, value=8, step=1)

with cfg_col5:
    show_markdown = st.checkbox("Exibir Markdown gerado", value=False)

size_row1, size_row2 = st.columns(2)
with size_row1:
    first_col_pct = st.slider("Largura 1ª coluna (%)", min_value=20, max_value=40, value=30, step=1)
with size_row2:
    pass

color_col1, color_col2, color_col3 = st.columns([1, 1, 2])

with color_col1:
    font_color_hex = st.color_picker("Cor da fonte PPTX", "#000000")

with color_col2:
    zebra_color_hex = st.color_picker("Cor zebra PPTX/Tabelas", "#FFFFFF")  # white per screenshot

with color_col3:
    zebra_alpha = st.slider(
        "Transparência zebra (PPTX/Tabelas)",
        min_value=0.0,
        max_value=1.0,
        value=1.0,  # fully opaque
        step=0.05,
    )

title_pos_col1, title_pos_col2 = st.columns(2)

with title_pos_col1:
    title_top_offset = st.slider("Posição título (in)", min_value=0.1, max_value=1.0, value=0.5, step=0.1)

with title_pos_col2:
    title_padding = st.slider(
        "Padding título/tabela (in)",
        min_value=0.3,
        max_value=2.0,
        value=0.7,  # per screenshot
        step=0.1,
    )

bg_image = st.file_uploader("Upload imagem de fundo (opcional)", type=["png", "jpg", "jpeg"])
bg_source = None
if bg_image:
    bg_source = bg_image

st.divider()

# ========== Add Page 1 ==========
st.subheader("Adicionar página 1 (Tabela Resumida)")
c1, c2 = st.columns([1, 3])

with c1:
    page1_name = st.text_input("Nome da Página 1 (título do slide)", value="Página 1")

with c2:
    placeholder_1 = (
        "Cole aqui o texto da Tabela Resumida (mesma estrutura da imagem 1). "
        "Apenas os percentuais variam.\n\n"
        "Asset Allocation\t\tConservador\t\t\tModerado\t\t\tSofisticado\n"
        "\t\tAlvo (%)\tFaixa Tolerada (%)\t\tAlvo (%)\tFaixa Tolerada (%)\t\tAlvo (%)\tFaixa Tolerada (%)\n"
        "Brasil\t\t95,0%\t90% - 100%\t\t87,5%\t80% - 100%\t\t82,0%\t60% - 100%\n"
        "Renda Fixa\t\t80,0%\t65% - 90%\t\t65,0%\t65% - 90%\t\t49,0%\t65% - 90%\n"
        "Pós Fixado\t\t60,0%\t45% - 65%\t\t45,0%\t35% - 55%\t\t34,0%\t30% - 50%\n"
        "Pré-Fixado\t\t0,0%\t0% - 10%\t\t0,0%\t0% - 10%\t\t0,0%\t0% - 10%\n"
        "Inflação\t\t20,0%\t15% - 25%\t\t20,0%\t15% - 25%\t\t15,0%\t10% - 20%\n"
        "Multimercado\t\t15,0%\t0% - 20%\t\t17,5%\t0% - 20%\t\t12,5%\t0% - 30%\n"
        "No Brasil\t\t10,0%\t0% - 15%\t\t12,5%\t0% - 15%\t\t12,5%\t0% - 20%\n"
        "No Exterior\t\t5,0%\t0% - 10%\t\t5,0%\t0% - 10%\t\t0,0%\t0% - 10%\n"
        "Renda Variável\t\t0,0%\t0% - 10%\t\t5,0%\t0% - 10%\t\t13,0%\t5% - 20%\n"
        "Ações\t\t0,0%\t0% - 5%\t\t0,0%\t0% - 5%\t\t5,0%\t0% - 15%\n"
        "Fundo de Ações\t\t0,0%\t0% - 5%\t\t0,0%\t0% - 5%\t\t3,0%\t0% - 15%\n"
        "FIIs | FIAgro | FIInfra\t\t0,0%\t0% - 5%\t\t5,0%\t0% - 5%\t\t5,0%\t0% - 15%\n"
        "Alternativos\t\t0,0%\t0% - 2%\t\t0,0%\t0% - 5%\t\t7,5%\t0% - 20%\n"
        "Global\t\t5,0%\t0% - 10%\t\t12,5%\t0% - 20%\t\t18,0%\t0% - 40%\n"
        "Renda Fixa\t\t5,0%\t0% - 10%\t\t7,5%\t0% - 15%\t\t8,0%\t0% - 25%\n"
        "Renda Variável\t\t0,0%\t0% - 5%\t\t5,0%\t0% - 10%\t\t10,0%\t0% - 25%\n"
        "TOTAL\t\t100,0%\t\t\t100,0%\t\t\t100,0%\n"
    )
    pasted_text1 = st.text_area(
        "Texto da Tabela (Resumida)",
        height=320,
        placeholder=placeholder_1,
        key="ta_page1",
    )

col_btn1 = st.columns([1, 1, 6])

with col_btn1[0]:
    if st.button("Adicionar Página 1", type="primary", key="btn_add_p1"):
        if not pasted_text1.strip():
            st.warning("Cole o texto contendo a tabela da Página 1.")
        else:
            try:
                df_simple, indents = parse_single_table_text(pasted_text1)
                df_view = to_multiindex_view(df_simple)
                name_final = (page1_name or "").strip() or "Página 1"
                st.session_state.pages.append(
                    {
                        "nome": name_final,
                        "df": df_view,
                        "indent": indents,
                        "subheader_extra": [],
                        "top_blacklist": [],
                    }
                )
                st.success(f"Página '{name_final}' adicionada.")
            except Exception as e:
                st.error("Erro ao parsear/reformatar a tabela da Página 1.")
                st.exception(e)

with col_btn1[1]:
    if st.button("Limpar lista", key="btn_clear"):
        st.session_state.pages = []
        st.info("Lista limpa.")

st.divider()

# ========== Add Page 2 ==========
st.subheader("Adicionar página 2 (Tabela Detalhada)")
c21, c22 = st.columns([1, 3])

with c21:
    page2_name = st.text_input("Nome da Página 2 (título do slide)", value="Página 2")

with c22:
    placeholder_2 = (
        "Cole aqui o texto da Tabela Detalhada (como no exemplo fornecido).\n\n"
        "Asset Allocation\t\tConservador\t\t\tModerado\t\t\tSofisticado\t\n"
        "\t\tAlvo (%)\tFaixa Tolerada (%)\t\tAlvo (%)\tFaixa Tolerada (%)\t\tAlvo (%)\tFaixa Tolerada (%)\n"
        "Liquidez\t\t17,5%\t10% - 50%\t\t12,5%\t10% - 40%\t\t10,0%\t0% - 30%\n"
        "Crédito Bancário - Liq Diaria\t\t0,0%\t0% - 25%\t\t0,0%\t0% - 20%\t\t0,0%\t0% - 15%\n"
        "Fundos Liquidez D+0 e D+1\t\t17,5%\t10% - 25%\t\t12,5%\t15% - 25%\t\t10,0%\t0% - 20%\n"
        "Renda Fixa - Títulos Públicos\t\t0,0%\t0% - 25%\t\t0,0%\t0% - 25%\t\t0,0%\t0% - 25%\n"
        "LFT- Selic\t\t0,0%\t0% - 10%\t\t0,0%\t0% - 10%\t\t0,0%\t0% - 10%\n"
        "LTN ou NTN-F - Pré\t\t0,0%\t0% - 10%\t\t0,0%\t0% - 10%\t\t0,0%\t0% - 10%\n"
        "NTN-B - IPCA+\t\t0,0%\t0% - 20%\t\t0,0%\t0% - 20%\t\t0,0%\t0% - 20%\n"
        "Renda Fixa - Crédito Bancário\t\t15,0%\t10% - 45%\t\t12,5%\t10% - 40%\t\t10,0%\t10% - 30%\n"
        "Bancários Pós CDI\t\t15,0%\t10% - 25%\t\t12,5%\t10% - 25%\t\t10,0%\t10% - 25%\n"
        "Bancários Pré\t\t0,0%\t0% - 15%\t\t0,0%\t0% - 15%\t\t0,0%\t0% - 15%\n"
        "Bancários IPCA+\t\t0,0%\t0% - 15%\t\t0,0%\t0% - 15%\t\t0,0%\t0% - 15%\n"
        "Renda Fixa - Crédito Privado\t\t25,0%\t0% - 40%\t\t25,0%\t0% - 40%\t\t22,0%\t0% - 40%\n"
        "Crédito % CDI ou CDI+\t\t5,0%\t0% - 15%\t\t5,0%\t0% - 15%\t\t7,0%\t0% - 15%\n"
        "Crédito Pré\t\t0,0%\t0% - 15%\t\t0,0%\t0% - 15%\t\t0,0%\t0% - 15%\n"
        "Crédito IPCA+\t\t20,0%\t10% - 25%\t\t20,0%\t10% - 25%\t\t15,0%\t10% - 25%\n"
        "Fundos de Renda Fixa\t\t22,5%\t10% - 40%\t\t20,0%\t10% - 35%\t\t20,0%\t10% - 30%\n"
        "Fundo de Crédito Isento\t\t15,0%\t10% - 25%\t\t15,0%\t10% - 25%\t\t15,0%\t10% - 25%\n"
        "Fundo de Crédito\t\t0,0%\t0% - 15%\t\t0,0%\t0% - 15%\t\t0,0%\t0% - 15%\n"
        "Previdência Privada - Renda Fixa\t\t7,5%\t0% - 15%\t\t5,0%\t0% - 15%\t\t5,0%\t0% - 15%\n"
        "Multimercados\t\t15,0%\t0% - 25%\t\t17,5%\t0% - 30%\t\t12,5%\t0% - 35%\n"
        "Multimercado Brasil\t\t5,0%\t0% - 15%\t\t7,5%\t0% - 15%\t\t7,5%\t0% - 25%\n"
        "Multimercado no Exterior\t\t5,0%\t0% - 15%\t\t5,0%\t0% - 15%\t\t0,0%\t0% - 15%\n"
        "Previdência Privada - Multimercado\t\t5,0%\t0% - 15%\t\t5,0%\t0% - 15%\t\t5,0%\t0% - 15%\n"
        "Renda Variável\t\t0,0%\t0% - 10%\t\t5,0%\t0% - 10%\t\t13,0%\t5% - 20%\n"
        "Ações\t\t0,0%\t0% - 5%\t\t0,0%\t0% - 5%\t\t5,0%\t0% - 15%\n"
        "Fundo de Ações\t\t0,0%\t0% - 5%\t\t0,0%\t0% - 5%\t\t3,0%\t0% - 15%\n"
        "FIIs | FIAgro | FIInfra\t\t0,0%\t0% - 5%\t\t5,0%\t0% - 15%\t\t5,0%\t0% - 20%\n"
        "Alternativos\t\t0,0%\t0% - 5%\t\t0,0%\t0% - 5%\t\t7,5%\t5% - 20%\n"
        "Criptomoedas\t\t0,0%\t0% - 5%\t\t0,0%\t0% - 5%\t\t2,5%\t0% - 10%\n"
        "Produtos Estruturados (FIDCs, FIPs, PEs)\t\t0,0%\t0% - 5%\t\t0,0%\t0% - 5%\t\t5,0%\t0% - 15%\n"
        "Global\t\t5,0%\t0% - 20%\t\t12,5%\t0% - 20%\t\t18,0%\t0% - 40%\n"
        "Global Renda Fixa Bonds\t\t0,0%\t0% - 5%\t\t0,0%\t0% - 10%\t\t0,0%\t0% - 10%\n"
        "Global Renda Fixa Fundos\t\t5,0%\t0% - 10%\t\t7,5%\t0% - 10%\t\t8,0%\t0% - 15%\n"
        "Global Renda Variável Ações\t\t0,0%\t0% - 5%\t\t0,0%\t0% - 10%\t\t5,0%\t0% - 15%\n"
        "Global Renda Variável Fundos / ETFs\t\t0,0%\t0% - 5%\t\t5,0%\t0% - 10%\t\t5,0%\t0% - 15%\n"
        "TOTAL\t\t100,0%\t\t\t100,0%\t\t\t100,0%\n"
    )
    pasted_text2 = st.text_area(
        "Texto da Tabela (Detalhada)",
        height=400,
        placeholder=placeholder_2,
        key="ta_page2",
    )

col_btn2 = st.columns([1, 6])

with col_btn2[0]:
    if st.button("Adicionar Página 2", type="primary", key="btn_add_p2"):
        if not pasted_text2.strip():
            st.warning("Cole o texto contendo a tabela da Página 2.")
        else:
            try:
                df_simple2, indents2 = parse_free_table_text(pasted_text2)
                df_view2 = to_multiindex_view(df_simple2)
                name_final2 = (page2_name or "").strip() or "Página 2"
                st.session_state.pages.append(
                    {
                        "nome": name_final2,
                        "df": df_view2,
                        "indent": indents2,
                        # Apply subheader style to Liquidez and Global on page 2
                        # and prevent Global from getting the dark-brown top bar.
                        "subheader_extra": ["Liquidez", "Global"],
                        "top_blacklist": ["Global"],
                    }
                )
                st.success(f"Página '{name_final2}' adicionada.")
            except Exception as e:
                st.error("Erro ao parsear/reformatar a tabela da Página 2.")
                st.exception(e)

st.divider()

# ========== Display Pages ==========
st.subheader("Páginas adicionadas")
if not st.session_state.pages:
    st.info("Nenhuma página adicionada ainda.")
else:
    for idx, p in enumerate(st.session_state.pages):
        with st.expander(f"{idx + 1}. {p['nome']}", expanded=False):
            st.markdown("**Tabela (layout estruturado)**")
            df = p["df"]
            indents = p["indent"]
            subextra = p.get("subheader_extra", [])
            topblack = p.get("top_blacklist", [])
            if not df.empty:
                st.table(
                    style_table(
                        df,
                        zebra_color_hex,
                        zebra_alpha,
                        indents,
                        first_col_pct,
                        subheader_extra=subextra,
                        top_blacklist=topblack,
                    )
                )
                if show_markdown:
                    st.caption("Markdown (Tabela)")
                    st.code(df_to_markdown(df))
            else:
                st.info("Tabela vazia ou não encontrada.")

            if st.button(f"Remover '{p['nome']}'", key=f"rm_{idx}"):
                st.session_state.pages.pop(idx)
                st.rerun()

st.divider()

# ========== Export ==========
if st.session_state.pages:
    try:
        ppt_bytes = build_pptx(
            nome_doc=nome_doc,
            pages=st.session_state.pages,
            title_font_size=title_font_size,
            hdr_size=hdr_font_size,
            body_size=body_font_size,
            title_top_offset=title_top_offset,
            title_padding=title_padding,
            bg_source=bg_source,
            font_color_hex=font_color_hex,
            zebra_color_hex=zebra_color_hex,
            zebra_alpha=zebra_alpha,
            first_col_pct=first_col_pct,
        )
        st.download_button(
            label="Baixar PowerPoint (.pptx)",
            data=ppt_bytes,
            file_name="asset_allocation_tabela_unica.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception as e:
        st.error("Erro ao gerar PPTX")
        st.exception(e)
else:
    st.info("Adicione ao menos uma página para exportar.")