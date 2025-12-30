import io
import os
import re
from typing import Dict, List, Optional, Tuple, Union

import pandas as pd
import streamlit as st
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

### Simple Authentication
if not st.session_state.get("authenticated", False):
    st.warning("Please enter the password on the Home page first.")
    st.write("Status: Não Autenticado")
    st.stop()

# prevent the rest of the page from running
st.write("Autenticado")

# -------------------- Constants --------------------

# Ceres Wealth style guide
FONT_NAME = "Century Gothic"

# Colors (Ceres Wealth palette)
TITLE_COLOR = "#825120"  # header background (web tables)
DEFAULT_ZEBRA_HEX = "#EFEFEF"  # default zebra color

# RGB tuples for PPTX
HEADER_BG_RGB = (0x82, 0x51, 0x20)  # #825120

# A4 Size (vertical)
SLIDE_WIDTH = Inches(8.27)
SLIDE_HEIGHT = Inches(11.69)

LIGHT_GRAY = RGBColor(0xEF, 0xEF, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CERES_BRONZE = RGBColor(0xB6, 0x84, 0x4C)

HDR_FONT_SIZE = 9
BODY_FONT_SIZE = 8
HDR_ROW_H_IN = 0.28
BODY_ROW_H_IN = 0.24

MARGIN_IN = 0.5


# -------------------- Utilities --------------------


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.strip().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def hex_to_rgb_tuple(hex_str: str) -> Tuple[int, int, int]:
    h = hex_str.strip().lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def hex_to_rgba_str(hex_str: str, alpha: float) -> str:
    h = hex_str.strip().lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    a = min(max(alpha, 0.0), 1.0)
    return f"rgba({r}, {g}, {b}, {a})"


def normalize_text(text: str) -> str:
    return (
        text.replace("\u00A0", " ")
        .replace("\r\n", "\n")
        .replace("\r", "\n")
        .strip()
    )


def sanitize_line(line: str) -> str:
    return re.sub(r"\s+", " ", line.strip())


def parts_count_with_content(line: str) -> int:
    """Count non-empty parts when split by tabs."""
    return len([p for p in line.split("\t") if p.strip()])


def extract_cnpj_pattern(text: str) -> Tuple[str, Optional[str]]:
    """
    Extract CNPJ pattern from text.
    Returns (text_without_cnpj, cnpj_part)
    CNPJ pattern: XX.XXX.XXX/XXXX-XX
    """
    cnpj_match = re.search(
        r"\b\d{1,2}\.\d{3}\.\d{3}/\d{4}-\d{2}\b", text
    )
    if cnpj_match:
        cnpj_part = cnpj_match.group(0)
        text_part = text[: cnpj_match.start()].rstrip()
        # Remove "CNPJ" text if it appears right before the CNPJ number
        text_part = re.sub(
            r"\s*CNPJ\s*$", "", text_part, flags=re.IGNORECASE
        ).rstrip()
        return (text_part, cnpj_part)
    return (text, None)


def extract_resumo(lines: List[str]) -> Dict[str, str]:
    """Extract the Resumo section metrics."""
    resumo = {}
    in_resumo = False
    sep_pattern = re.compile(r"\t+| {2,}")

    for i, ln in enumerate(lines):
        if sanitize_line(ln).lower() == "resumo":
            in_resumo = True
            continue

        if in_resumo:
            s = sanitize_line(ln)
            if not s:
                continue

            s_lower = s.lower()
            if ("categoria" in s_lower and "ativo" in s_lower) or (
                "rent" in s_lower and "%cdi" in s_lower
            ):
                break

            parts = [p.strip() for p in sep_pattern.split(ln) if p.strip()]
            if len(parts) >= 2:
                key = parts[0]
                value = parts[1]
                resumo[key] = value

    return resumo


def parse_new_portfolio_format(raw_text: str) -> Dict:
    """
    Parse the tab or multi-space separated portfolio format.
    Returns:
      {
        'titulo_guess': Optional[str],
        'composicao': List[Tuple[str, Optional[float], str]],
        'holdings': List[Dict[str, str]],
        'resumo': Dict[str, str]
      }
    """
    text = normalize_text(raw_text)
    lines = text.split("\n")

    # Extract title (first non-empty line not a known section)
    titulo_guess = None
    for ln in lines:
        s = sanitize_line(ln)
        if s and s.lower() not in {
            "composição da carteira",
            "alocação por classe",
            "alocacao por classe",
            "resumo",
        }:
            titulo_guess = s
            break

    # Extract Composição da Carteira / Alocação por Classe
    composicao = []
    sep_pattern = re.compile(r"\t+| {2,}")
    for i, ln in enumerate(lines):
        head = sanitize_line(ln).lower()
        if head in {
            "composição da carteira",
            "alocação por classe",
            "alocacao por classe",
        }:
            i += 1
            while i < len(lines):
                raw = lines[i]
                s = sanitize_line(raw)
                s_lower = s.lower()

                if not raw.strip():
                    i += 1
                    continue
                if s_lower == "resumo":
                    break

                parts = [
                    p.strip()
                    for p in sep_pattern.split(raw)
                    if p.strip()
                ]
                if len(parts) >= 2:
                    classe = parts[0]
                    pct_str = parts[1]
                    if pct_str and not pct_str.endswith("%"):
                        pct_str = pct_str + "%"
                    try:
                        pct_float = float(
                            pct_str.rstrip("%")
                            .replace(".", "")
                            .replace(",", ".")
                        )
                    except Exception:
                        pct_float = None
                    if classe:
                        composicao.append((classe, pct_float, pct_str))
                i += 1
            break

    # Extract Resumo
    resumo = extract_resumo(lines)

    # Extract Holdings - robust logic for
    # "Categoria  Ativo  %  Rent. 12M  %CDI"
    holdings: List[Dict[str, str]] = []
    current_section = None
    i = 0

    # detect header first
    header_regex = re.compile(r"categoria\s+ativo", re.IGNORECASE)
    while i < len(lines):
        ln = lines[i]
        s = sanitize_line(ln)
        s_lower = s.lower()
        if not s:
            i += 1
            continue
        if header_regex.search(s_lower) and (
            "%cdi" in s_lower or "rent" in s_lower or "12m" in s_lower
        ):
            i += 1
            break
        i += 1

    # parse rows after header
    while i < len(lines):
        ln = lines[i]
        if not ln.strip():
            i += 1
            continue

        parts = [
            p.strip()
            for p in re.split(r"\t+| {2,}", ln)
            if p.strip()
        ]
        # Expect 5 columns: Categoria, Ativo, %, Rent. 12M, %CDI
        if len(parts) >= 4:
            group = parts[0]
            ativo = parts[1]
            peso_str = parts[2] if len(parts) > 2 else ""
            ret12_str = parts[3] if len(parts) > 3 else ""
            cdi_str = parts[4] if len(parts) > 4 else ""

            if peso_str and not peso_str.endswith("%"):
                peso_str += "%"
            if ret12_str and not ret12_str.endswith("%"):
                ret12_str += "%"
            if cdi_str and not cdi_str.endswith("%"):
                cdi_str += "%"

            holdings.append(
                {
                    "grupo": group,
                    "secao": current_section or "Sem Seção",
                    "ativo": ativo,
                    "peso": peso_str,
                    "ret12m": ret12_str,
                    "cdi12m": cdi_str,
                }
            )
        else:
            s_lower = sanitize_line(ln).lower()
            if s_lower in {
                "resumo",
                "composição da carteira",
                "alocação por classe",
                "alocacao por classe",
            }:
                break
        i += 1

    return {
        "titulo_guess": titulo_guess,
        "composicao": composicao,
        "holdings": holdings,
        "resumo": resumo,
    }


# -------------------- PPTX helpers --------------------


def set_paragraph(
    paragraph,
    text: str,
    size: int,
    bold: bool = False,
    color: Optional[str] = None,
    align: Optional[int] = None,
):
    if hasattr(paragraph, "clear"):
        paragraph.clear()
    else:
        paragraph.text = ""
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = FONT_NAME
    font.size = Pt(size)
    font.bold = bold
    if color:
        font.color.rgb = hex_to_rgb(color)
    if align is not None:
        paragraph.alignment = align


def set_cell_fill_with_alpha(
    cell, rgb_tuple: Tuple[int, int, int], alpha: float
):
    """
    Set cell fill with alpha channel.
    rgb_tuple: (R, G, B)
    alpha: 0.0 (fully transparent) to 1.0 (fully opaque)
    """
    from pptx.oxml.xmlchemy import OxmlElement

    try:
        tcPr = cell._tc.get_or_add_tcPr()

        # Remove existing fills
        for fill in tcPr.findall(qn("a:solidFill")):
            tcPr.remove(fill)

        # Create new solidFill
        solidFill = OxmlElement("a:solidFill")

        srgbClr = OxmlElement("a:srgbClr")
        hex_val = "{:02X}{:02X}{:02X}".format(
            rgb_tuple[0], rgb_tuple[1], rgb_tuple[2]
        )
        srgbClr.set("val", hex_val)

        # Add alpha (100000 = opaque, 0 = transparent)
        alpha_elem = OxmlElement("a:alpha")
        alpha_elem.set("val", str(int(alpha * 100000)))
        srgbClr.append(alpha_elem)

        solidFill.append(srgbClr)
        tcPr.append(solidFill)
    except Exception:
        pass


def set_cell_text_centered(
    cell,
    text: str,
    size: int,
    bold: bool = False,
    font_color_hex: Optional[str] = None,
):
    """Set cell text with center alignment and vertical center."""
    cell.text_frame.clear()
    cell.text_frame.word_wrap = True
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    color_hex = font_color_hex or "#000000"
    run.font.color.rgb = hex_to_rgb(color_hex)


def set_cell_text_with_cnpj(
    cell,
    text: str,
    size: int,
    cnpj_size: int,
    bold: bool = False,
    font_color_hex: Optional[str] = None,
):
    """
    Set cell text, making CNPJ portion smaller font if present.
    Handles cases where text is only CNPJ or contains CNPJ.
    Always center-aligned.
    """
    text_part, cnpj_part = extract_cnpj_pattern(text)

    cell.text_frame.clear()
    cell.text_frame.word_wrap = True
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    color_hex = font_color_hex or "#000000"
    rgb = hex_to_rgb(color_hex)

    # Add main text (if it exists and is not just whitespace)
    if text_part and text_part.strip():
        run = p.add_run()
        run.text = text_part
        run.font.name = FONT_NAME
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb

    # Add CNPJ if present (always use smaller font)
    if cnpj_part:
        # Add newline before CNPJ if there was text before it
        if text_part and text_part.strip():
            run = p.add_run()
            run.text = "\n"
            run.font.size = Pt(size)

        run = p.add_run()
        run.text = cnpj_part
        run.font.name = FONT_NAME
        run.font.size = Pt(cnpj_size)
        run.font.bold = False
        run.font.color.rgb = rgb
    elif not text_part or not text_part.strip():
        # Fallback: add original text if nothing was extracted
        run = p.add_run()
        run.text = text
        run.font.name = FONT_NAME
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb


# --------- Borderless table theme ---------


def apply_borderless_theme_to_table(table) -> None:
    """
    Apply a borderless theme to the entire table at the XML level.
    This removes all cell borders in one operation by modifying the
    table's underlying XML structure.
    """
    from pptx.oxml.xmlchemy import OxmlElement

    try:
        tbl = table._tbl
        tblPr = tbl.tblPr

        # Iterate through all cells and remove border elements
        for tc in tbl.iter(qn("a:tc")):
            tcPr = tc.get_or_add_tcPr()

            # Remove all existing border line elements (lnL, lnR, lnT, lnB)
            for border_tag in [
                qn("a:lnL"),
                qn("a:lnR"),
                qn("a:lnT"),
                qn("a:lnB"),
            ]:
                for border_elem in tcPr.findall(border_tag):
                    tcPr.remove(border_elem)

            # Set all borders to "no fill" (invisible)
            for border_pos in ["lnL", "lnR", "lnT", "lnB"]:
                ln = OxmlElement(f"a:{border_pos}")
                # Add noFill to make the line invisible
                noFill = OxmlElement("a:noFill")
                ln.append(noFill)
                tcPr.append(ln)

    except Exception:
        pass


def add_background(
    slide,
    prs: Presentation,
    bg_source: Optional[Union[str, io.BytesIO]] = None,
):
    """Set slide background image (file path or file-like object)."""
    if bg_source is None:
        return
    try:
        slide.shapes.add_picture(
            bg_source,
            Inches(0),
            Inches(0),
            prs.slide_width,
            prs.slide_height,
        )
    except Exception:
        pass


def add_cover_slide(
    prs: Presentation,
    title: str,
    bg_source: Optional[Union[str, io.BytesIO]] = None,
    font_color_hex: str = "#000000",
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add background image covering the entire slide
    add_background(slide, prs, bg_source)

    left = Inches(MARGIN_IN)
    top = Inches(3.2)
    width = prs.slide_width - Inches(2 * MARGIN_IN)
    height = Inches(1.8)
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    # Cover title uses chosen font color
    set_paragraph(p, title, 44, True, font_color_hex, PP_ALIGN.CENTER)


def add_portfolio_slide(
    prs: Presentation,
    nome: str,
    parsed: Dict,
    title_font_size: int,
    hdr_size: int,
    body_size: int,
    cnpj_size: int,
    title_top_offset: float,
    title_padding: float,
    bg_source: Optional[Union[str, io.BytesIO]] = None,
    font_color_hex: str = "#000000",
    zebra_color_hex: str = DEFAULT_ZEBRA_HEX,
    zebra_alpha: float = 1.0,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add background image covering the entire slide
    add_background(slide, prs, bg_source)

    margin = Inches(MARGIN_IN)
    content_w = prs.slide_width - 2 * margin

    # Zebra color for PPTX
    zebra_rgb = hex_to_rgb_tuple(zebra_color_hex)

    # Title
    title_top = Inches(title_top_offset)
    tb_title = slide.shapes.add_textbox(
        margin, title_top, content_w, Inches(0.7)
    )
    p = tb_title.text_frame.paragraphs[0]
    set_paragraph(p, nome, title_font_size, True, font_color_hex)

    current_top = title_top + Inches(title_padding)

    # Get data
    composicao = parsed.get("composicao", [])
    resumo = parsed.get("resumo", {})
    rows = parsed.get("holdings", [])

    hdr_h = Inches(HDR_ROW_H_IN)
    row_h = Inches(BODY_ROW_H_IN)

    # ========== SIDE BY SIDE: Resumo (LEFT) + Composição (RIGHT) ==========
    if composicao or resumo:
        gap = Inches(0.15)
        content_w_emu = int(content_w)
        gap_emu = int(gap)
        half_w_emu = (content_w_emu - gap_emu) // 2
        half_w = Inches(half_w_emu / 914400)

        hdr_top = current_top - Inches(0.32)
        tbl_top = current_top

        # LEFT: Resumo
        if resumo:
            tb_hdr_resumo = slide.shapes.add_textbox(
                margin, hdr_top, half_w, Inches(0.3)
            )
            p = tb_hdr_resumo.text_frame.paragraphs[0]
            # Section subtitles use chosen font color
            set_paragraph(
                p,
                "Resumo",
                15,
                True,
                font_color_hex,
            )

            resumo_rows = 1 + len(resumo)
            resumo_table_h = hdr_h + row_h * (resumo_rows - 1)

            resumo_shape = slide.shapes.add_table(
                resumo_rows, 2, margin, tbl_top, half_w, resumo_table_h
            )
            resumo_table = resumo_shape.table

            # Apply borderless theme to table
            apply_borderless_theme_to_table(resumo_table)

            resumo_table.rows[0].height = hdr_h
            for i in range(1, resumo_rows):
                resumo_table.rows[i].height = row_h

            desc_w = int(half_w_emu * 0.6)
            val_w = int(half_w_emu * 0.4)
            resumo_table.columns[0].width = desc_w
            resumo_table.columns[1].width = val_w

            # Header with solid primary color, white font (style guide)
            for j, h in enumerate(["Métrica", "Valor"]):
                cell = resumo_table.cell(0, j)
                set_cell_fill_with_alpha(cell, HEADER_BG_RGB, 1.0)
                set_cell_text_centered(
                    cell, h, hdr_size, True, "#FFFFFF"
                )

            # Rows: zebra color / white, font color from selector
            for i, (key, val) in enumerate(resumo.items(), start=1):
                vals = [key, val]
                is_even = i % 2 == 0

                for j, v in enumerate(vals):
                    cell = resumo_table.cell(i, j)

                    if is_even:
                        set_cell_fill_with_alpha(
                            cell, zebra_rgb, zebra_alpha
                        )
                    else:
                        cell.fill.background()

                    set_cell_text_centered(
                        cell, v, body_size, False, font_color_hex
                    )

        # RIGHT: Composição
        if composicao:
            comp_left = margin + half_w + gap

            tb_hdr_comp = slide.shapes.add_textbox(
                comp_left, hdr_top, half_w, Inches(0.3)
            )
            p = tb_hdr_comp.text_frame.paragraphs[0]
            set_paragraph(
                p,
                "Alocação por Classe",
                15,
                True,
                font_color_hex,
            )

            comp_rows = 1 + len(composicao)
            comp_table_h = hdr_h + row_h * (comp_rows - 1)

            comp_shape = slide.shapes.add_table(
                comp_rows, 2, comp_left, tbl_top, half_w, comp_table_h
            )
            comp_table = comp_shape.table

            # Apply borderless theme to table
            apply_borderless_theme_to_table(comp_table)

            comp_table.rows[0].height = hdr_h
            for i in range(1, comp_rows):
                comp_table.rows[i].height = row_h

            peso_w = int(half_w_emu * 0.35)
            classe_w = int(half_w_emu * 0.65)
            comp_table.columns[0].width = classe_w
            comp_table.columns[1].width = peso_w

            # Header with solid primary color, white font
            for j, h in enumerate(["Classe", "Peso"]):
                cell = comp_table.cell(0, j)
                set_cell_fill_with_alpha(cell, HEADER_BG_RGB, 1.0)
                set_cell_text_centered(
                    cell, h, hdr_size, True, "#FFFFFF"
                )

            # Rows: zebra color / white - ALL CENTER ALIGNED
            for i, (classe, _pct_float, pct_str) in enumerate(
                composicao, start=1
            ):
                vals = [classe, pct_str]
                is_even = i % 2 == 0

                for j, val in enumerate(vals):
                    cell = comp_table.cell(i, j)

                    if is_even:
                        set_cell_fill_with_alpha(
                            cell, zebra_rgb, zebra_alpha
                        )
                    else:
                        cell.fill.background()

                    set_cell_text_centered(
                        cell, val, body_size, False, font_color_hex
                    )

        # Calculate max height of both side-by-side tables
        max_resumo_height = (
            (hdr_h + row_h * len(resumo)) if resumo else Inches(0)
        )
        max_comp_height = (
            (hdr_h + row_h * len(composicao))
            if composicao
            else Inches(0)
        )
        max_side_by_side_height = max(
            max_resumo_height, max_comp_height
        )

        # Update current_top for full-width Ativos below
        current_top = tbl_top + max_side_by_side_height + Inches(0.4)

    # ========== FULL WIDTH: Ativos ==========
    if rows:
        # Calculate the actual width used by the side-by-side tables
        actual_content_w = 2.1 * half_w + gap

        tb_hdr_ativos = slide.shapes.add_textbox(
            margin,
            current_top - Inches(0.32),
            actual_content_w,
            Inches(0.3),
        )
        p = tb_hdr_ativos.text_frame.paragraphs[0]
        set_paragraph(
            p,
            "Alocação",
            15,
            True,
            font_color_hex,
        )

        pos_rows = 1 + len(rows)
        pos_table_h = hdr_h + row_h * (pos_rows - 1)

        pos_shape = slide.shapes.add_table(
            pos_rows,
            5,
            margin,
            current_top,
            actual_content_w,
            pos_table_h,
        )
        pos_table = pos_shape.table

        # Apply borderless theme to table
        apply_borderless_theme_to_table(pos_table)

        pos_table.rows[0].height = hdr_h
        for i in range(1, pos_rows):
            pos_table.rows[i].height = row_h

        # Column widths for FULL WIDTH
        actual_content_w_emu = int(actual_content_w)
        col_widths = [
            int(0.21 * actual_content_w_emu),  # grupo
            int(0.35 * actual_content_w_emu),  # ativo
            int(0.13 * actual_content_w_emu),  # peso
            int(0.13 * actual_content_w_emu),  # ret12m
            int(0.13 * actual_content_w_emu),  # cdi12m
        ]
        for j, w in enumerate(col_widths):
            pos_table.columns[j].width = w

        # Header with solid primary color, white font
        headers = ["Grupo", "Ativo", "%", "Rent. 12M", "%CDI 12M"]
        for j, h in enumerate(headers):
            cell = pos_table.cell(0, j)
            set_cell_fill_with_alpha(cell, HEADER_BG_RGB, 1.0)
            set_cell_text_centered(
                cell, h, hdr_size, True, "#FFFFFF"
            )

        # Data rows - zebra color / white - ALL CENTER ALIGNED
        for i in range(1, pos_rows):
            if i <= len(rows):
                r = rows[i - 1]
                vals = [
                    r["grupo"],
                    r["ativo"],
                    r["peso"],
                    r["ret12m"],
                    r["cdi12m"],
                ]
            else:
                vals = ["", "", "", "", ""]

            is_even = i % 2 == 0

            for j, val in enumerate(vals):
                cell = pos_table.cell(i, j)

                if is_even:
                    set_cell_fill_with_alpha(
                        cell, zebra_rgb, zebra_alpha
                    )
                else:
                    cell.fill.background()

                if j == 1 and val:
                    # Ativo column with possible CNPJ
                    set_cell_text_with_cnpj(
                        cell,
                        val,
                        body_size,
                        cnpj_size,
                        False,
                        font_color_hex,
                    )
                else:
                    set_cell_text_centered(
                        cell, val, body_size, False, font_color_hex
                    )


def build_pptx(
    nome_doc: str,
    carteiras: List[Dict],
    title_font_size: int,
    hdr_size: int,
    body_size: int,
    cnpj_size: int,
    title_top_offset: float,
    title_padding: float,
    bg_source: Optional[Union[str, io.BytesIO]] = None,
    font_color_hex: str = "#000000",
    zebra_color_hex: str = DEFAULT_ZEBRA_HEX,
    zebra_alpha: float = 1.0,
) -> io.BytesIO:
    prs = Presentation()

    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    if nome_doc.strip():
        add_cover_slide(
            prs, nome_doc.strip(), bg_source, font_color_hex
        )

    for c in carteiras:
        add_portfolio_slide(
            prs,
            c["nome"],
            c["parsed"],
            title_font_size,
            hdr_size,
            body_size,
            cnpj_size,
            title_top_offset,
            title_padding,
            bg_source,
            font_color_hex,
            zebra_color_hex,
            zebra_alpha,
        )

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


# -------------------- Streamlit helpers --------------------


def style_table(
    df: pd.DataFrame, zebra_hex: str, zebra_alpha: float
):
    zebra_bg_css = hex_to_rgba_str(zebra_hex, zebra_alpha)

    styler = df.style.set_table_styles(
        [
            {
                "selector": "th",
                "props": [
                    ("background-color", TITLE_COLOR),
                    ("color", "white"),
                    ("font-weight", "bold"),
                    ("font-family", "'Century Gothic', sans-serif"),
                    ("border", "none"),
                ],
            },
            {
                "selector": "td",
                "props": [
                    ("font-family", "'Century Gothic', sans-serif"),
                    ("border", "none"),
                ],
            },
            {
                "selector": "table",
                "props": [
                    ("border-collapse", "collapse"),
                    ("border", "none"),
                ],
            },
        ]
    )

    def zebra(row):
        color = zebra_bg_css if row.name % 2 == 0 else "#FFFFFF"
        return [f"background-color: {color}"] * len(row)

    return styler.apply(zebra, axis=1)


# -------------------- Streamlit UI --------------------

st.set_page_config(
    page_title="Portfólios (Novo Formato) + PPTX", layout="centered"
)

st.markdown(
    "<h2 style='text-align: center;'>Parser de Portfólios (Novo Formato Tab) + PPTX</h2>",
    unsafe_allow_html=True,
)

st.markdown("<br><br><br>", unsafe_allow_html=True)

st.markdown(
    "Cole o texto no novo formato com abas (tab-separated) ou com 2+ "
    "espaços.\nSuporta: Composição/Alocação, Resumo e ativos com cabeçalho "
    "'Categoria Ativo % Rent. 12M %CDI'."
)

if "carteiras" not in st.session_state:
    st.session_state.carteiras = []

# ========== Configuration Section ==========
st.subheader("Configurações")

cfg_col1, cfg_col2, cfg_col3, cfg_col4, cfg_col5 = st.columns(5)

with cfg_col1:
    nome_doc = st.text_input("Nome do Documento (capa PPT)", value="")

with cfg_col2:
    # Guia: título principal 14
    title_font_size = st.slider(
        "Tamanho título PPT",
        min_value=10,
        max_value=36,
        value=13,
        step=1,
    )

with cfg_col3:
    # Guia: subtítulos 10
    hdr_font_size = st.slider(
        "Tamanho fonte headers",
        min_value=8,
        max_value=16,
        value=9,
        step=1,
    )

with cfg_col4:
    # Guia: texto 10
    body_font_size = st.slider(
        "Tamanho fonte corpo",
        min_value=7,
        max_value=14,
        value=8,
        step=1,
    )

with cfg_col5:
    cnpj_font_size = st.slider(
        "Tamanho fonte CNPJ",
        min_value=5,
        max_value=10,
        value=6,
        step=1,
    )

# Color configuration (font + zebra)
color_col1, color_col2, color_col3 = st.columns([1, 1, 2])

with color_col1:
    font_color_hex = st.color_picker(
        "Cor da fonte PPTX", "#000000"
    )

with color_col2:
    zebra_color_hex = st.color_picker(
        "Cor zebra PPTX/Tabelas", DEFAULT_ZEBRA_HEX
    )

with color_col3:
    zebra_alpha = st.slider(
        "Transparência zebra (PPTX/Tabelas)",
        min_value=0.0,
        max_value=1.0,
        value=1.0,
        step=0.05,
    )

title_pos_col1, title_pos_col2, title_pos_col3, title_pos_col4 = st.columns(
    4
)

with title_pos_col1:
    title_top_offset = st.slider(
        "Posição título (in)",
        min_value=0.1,
        max_value=1.0,
        value=0.5,
        step=0.1,
    )

with title_pos_col2:
    title_padding = st.slider(
        "Padding título/tabelas (in)",
        min_value=0.3,
        max_value=2.0,
        value=1.0,
        step=0.1,
    )

bg_image = st.file_uploader(
    "Upload imagem de fundo (opcional)", type=["png", "jpg", "jpeg"]
)
bg_source = None
if bg_image:
    bg_source = bg_image

st.divider()

# ========== Add Portfolio Section ==========
st.subheader("Adicionar portfólio")
c1, c2 = st.columns([1, 3])
with c1:
    nome_port = st.text_input("Nome do Portfólio (opcional)")
with c2:
    placeholder = (
        "Cole aqui o texto em formato tab-separated ou com 2+ espaços...\n"
        "Exemplo:\n"
        "ASSESSORIA DIGITAL - CONSERVADOR\n"
        "Alocação por Classe\n"
        "Pos Fixado\t50,00%\n"
        "...\n"
        "Resumo\n"
        "Valor da carteira\tR$ 300.000,00\n"
        "...\n"
        "Categoria\tAtivo\t%\tRent. 12M\t%CDI\n"
        "..."
    )
    texto = st.text_area(
        "Texto do Portfólio", height=260, placeholder=placeholder
    )

col_btn = st.columns([1, 1, 6])
with col_btn[0]:
    if st.button("Adicionar à lista", type="primary"):
        if not texto.strip():
            st.warning("Cole o texto do portfólio.")
            # keep font/zebra selections in state automatically via Streamlit
        else:
            try:
                parsed = parse_new_portfolio_format(texto)
                nome_final = (
                    (nome_port or "").strip()
                    or parsed.get("titulo_guess")
                    or "Portfólio"
                )
                st.session_state.carteiras.append(
                    {"nome": nome_final, "parsed": parsed}
                )
                st.success(f"Portfólio '{nome_final}' adicionado.")
            except Exception as e:
                st.error("Erro ao parsear o portfólio.")
                st.exception(e)

with col_btn[1]:
    if st.button("Limpar lista"):
        st.session_state.carteiras = []
        st.info("Lista limpa.")

st.divider()

# ========== Display Portfolios Section ==========
st.subheader("Portfólios adicionados")
if not st.session_state.carteiras:
    st.info("Nenhum portfólio adicionado ainda.")
else:
    for idx, c in enumerate(st.session_state.carteiras):
        parsed = c["parsed"]
        with st.expander(f"{idx + 1}. {c['nome']}", expanded=False):
            # Composição
            st.markdown("**Composição da Carteira**")
            comp = parsed.get("composicao", [])
            if comp:
                df_comp = pd.DataFrame(
                    [(a, c_str) for a, _c, c_str in comp],
                    columns=["Classe", "Peso (%)"],
                )
                st.table(
                    style_table(
                        df_comp,
                        zebra_color_hex,
                        zebra_alpha,
                    )
                )
            else:
                st.info("Sem composição.")

            # Resumo
            resumo = parsed.get("resumo", {})
            if resumo:
                st.markdown("**Resumo**")
                df_resumo = pd.DataFrame(
                    list(resumo.items()),
                    columns=["Métrica", "Valor"],
                )
                st.table(
                    style_table(
                        df_resumo,
                        zebra_color_hex,
                        zebra_alpha,
                    )
                )

            # Holdings
            st.markdown("**Ativos da carteira**")
            rows = parsed.get("holdings", [])
            if rows:
                df_h = (
                    pd.DataFrame(
                        rows,
                        columns=[
                            "grupo",
                            "ativo",
                            "peso",
                            "ret12m",
                            "cdi12m",
                        ],
                    ).rename(
                        columns={
                            "grupo": "Grupo",
                            "ativo": "Ativo",
                            "peso": "Peso (%)",
                            "ret12m": "Retorno 12M",
                            "cdi12m": "%CDI 12M",
                        }
                    )
                )
                st.dataframe(
                    df_h, width='stretch', hide_index=True
                )
            else:
                st.warning("Nenhum ativo identificado.")

            if st.button(f"Remover '{c['nome']}'", key=f"rm_{idx}"):
                st.session_state.carteiras.pop(idx)
                st.rerun()

st.divider()

# ========== Export Section ==========
if st.session_state.carteiras:
    try:
        ppt_bytes = build_pptx(
            nome_doc,
            st.session_state.carteiras,
            title_font_size,
            hdr_font_size,
            body_font_size,
            cnpj_font_size,
            title_top_offset,
            title_padding,
            bg_source,
            font_color_hex,
            zebra_color_hex,
            zebra_alpha,
        )
        st.download_button(
            label="Baixar PowerPoint (.pptx)",
            data=ppt_bytes,
            file_name="portfolios.pptx",
            mime=(
                "application/vnd.openxmlformats-officedocument."
                "presentationml.presentation"
            ),
        )
    except Exception as e:
        st.error("Erro ao gerar PPTX")
        st.exception(e)
else:
    st.info("Adicione pelo menos um portfólio para exportar.")