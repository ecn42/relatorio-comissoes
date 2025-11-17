import io
import os
import re
from decimal import Decimal, InvalidOperation
from typing import Dict, List, Optional, Tuple, Union

import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

### Simple Authentication
if not st.session_state.get("authenticated", False):
    st.warning("Please enter the password on the Home page first.")
    st.write("Status: Não Autenticado")
    st.stop()
    
  # prevent the rest of the page from running
st.write("Autenticado")





# -------------------- Constants --------------------

PALETTE_HEX = [
    "#013220",
    "#8c6239",
    "#57575a",
    "#b08568",
    "#09202e",
    "#582308",
    "#7a6200",
]

TITLE_COLOR = "#8C6239"
FONT_NAME = "Montserrat"  # must be installed on your machine for PPT

LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# PPT sizing
HDR_FONT_SIZE = 11  # table header
BODY_FONT_SIZE = 10  # table body
HDR_ROW_H_IN = 0.28  # inches
BODY_ROW_H_IN = 0.24  # inches

LEFT_COL_W_IN = 5.3  # left column width (Descrição, Valores, Setores)
GAP_IN = 0.25

# -------------------- Utilities --------------------


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.strip().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def normalize_newlines_and_spaces(text: str) -> str:
    return (
        text.replace("\u00A0", " ")
        .replace("\r\n", "\n")
        .replace("\r", "\n")
        .strip()
    )


def sanitize_line(line: str) -> str:
    return re.sub(r"\s+", " ", line.strip())


def find_exact_header(lines: List[str], header: str) -> int:
    for i, line in enumerate(lines):
        if sanitize_line(line).lower() == header.lower():
            return i
    return -1


def find_startswith_line(lines: List[str], startswith: str) -> int:
    sw = startswith.lower()
    for i, line in enumerate(lines):
        if sanitize_line(line).lower().startswith(sw):
            return i
    return -1


def next_non_empty_index(lines: List[str], start_idx: int) -> Optional[int]:
    i = start_idx + 1
    while i < len(lines):
        if lines[i].strip():
            return i
        i += 1
    return None


def parse_brl_amount(amount_str: str) -> Optional[Decimal]:
    try:
        s = amount_str.strip()
        s = s.replace("R$", "").strip()
        s = s.replace(".", "").replace(",", ".")
        return Decimal(s)
    except (InvalidOperation, AttributeError):
        return None


def format_brl(amount: Decimal) -> str:
    q = amount.quantize(Decimal("0.01"))
    whole, _, frac = f"{q:.2f}".partition(".")
    parts = []
    while len(whole) > 3:
        parts.append(whole[-3:])
        whole = whole[:-3]
    if whole:
        parts.append(whole)
    parts.reverse()
    return f"R$ {'.'.join(parts)},{frac}"


def is_dash_line(s: str) -> bool:
    s = s.strip()
    return bool(s) and all(ch in "-–—" for ch in s)


# -------------------- Parsing --------------------


def parse_sections(raw_text: str) -> Dict:
    text = normalize_newlines_and_spaces(raw_text)
    lines = text.split("\n")

    HDR_DESC = "Descrição"
    HDR_PORT = "Composição da carteira"
    HDR_DIST = "Distribuição por setor"

    LBL_APORTE = "Valor mínimo de aporte:"
    LBL_PERMAN = "Valor mínimo de permanência:"
    LBL_MOV = "Valor mínimo de movimentação:"

    idx_desc = find_exact_header(lines, HDR_DESC)
    idx_port = find_exact_header(lines, HDR_PORT)
    idx_dist = find_exact_header(lines, HDR_DIST)

    idx_aporte = find_startswith_line(lines, LBL_APORTE)
    idx_perman = find_startswith_line(lines, LBL_PERMAN)
    idx_mov = find_startswith_line(lines, LBL_MOV)

    descricao = None
    if idx_desc != -1:
        candidates = [
            i
            for i in [idx_aporte, idx_perman, idx_mov, idx_port, idx_dist]
            if i != -1 and i > idx_desc
        ]
        end = min(candidates) if candidates else len(lines)
        block = [l.strip() for l in lines[idx_desc + 1 : end] if l.strip()]
        if block:
            descricao = " ".join(block)

    def read_val_after(idx_label: int) -> Optional[str]:
        if idx_label == -1:
            return None
        idx_val = next_non_empty_index(lines, idx_label)
        if idx_val is None:
            return None
        return lines[idx_val].strip()

    aporte_raw = read_val_after(idx_aporte)
    perman_raw = read_val_after(idx_perman)
    mov_raw = read_val_after(idx_mov)

    def normalize_brl(raw_val: Optional[str]) -> Optional[str]:
        if not raw_val:
            return None
        amt = parse_brl_amount(raw_val)
        if amt is None:
            return raw_val.strip()
        return format_brl(amt)

    aporte_fmt = normalize_brl(aporte_raw)
    perman_fmt = normalize_brl(perman_raw)
    mov_fmt = normalize_brl(mov_raw)

    carteira_rows: List[Tuple[str, str, int]] = []
    if idx_port != -1:
        end = idx_dist if idx_dist != -1 else len(lines)
        table_lines = [l for l in lines[idx_port + 1 : end] if l.strip()]
        for ln in table_lines:
            s = sanitize_line(ln).lower()
            if "empresa" in s and "ticker" in s and "peso" in s:
                continue
            toks = ln.strip().split()
            if len(toks) < 3:
                continue
            weight_str = toks[-1]
            if not re.fullmatch(r"\d{1,3}\s*%", weight_str):
                continue
            ticker = toks[-2]
            empresa = " ".join(toks[:-2]).strip()
            try:
                peso = int(re.sub(r"\D", "", weight_str))
            except ValueError:
                continue
            if empresa and ticker and 0 <= peso <= 100:
                carteira_rows.append((empresa, ticker, peso))

    setores: List[Tuple[str, int]] = []
    if idx_dist != -1:
        i = idx_dist + 1
        while i < len(lines):
            ln = lines[i].strip()
            if not ln:
                i += 1
                continue
            m = re.fullmatch(r"(\d{1,3})\s*%", ln)
            if m:
                pct = int(m.group(1))
                j = i + 1
                cat = None
                while j < len(lines):
                    cand = lines[j].strip()
                    if not cand or is_dash_line(cand):
                        j += 1
                        continue
                    if re.fullmatch(r"\d{1,3}\s*%", cand):
                        break
                    cat = sanitize_line(cand)
                    break
                if cat is not None:
                    setores.append((cat, pct))
                    i = j + 1
                    continue
                else:
                    i += 1
                    continue
            else:
                i += 1

    return {
        "descricao": descricao,
        "aporte": aporte_fmt,
        "perman": perman_fmt,
        "mov": mov_fmt,
        "carteira": carteira_rows,
        "setores": setores,
    }


# -------------------- PPTX helpers --------------------


def set_paragraph(
    paragraph,
    text: str,
    size: int,
    bold: bool = False,
    color: Optional[str] = None,
    align: Optional[int] = None,
):
    if hasattr(paragraph, "clear"):
        paragraph.clear()
    else:
        paragraph.text = ""
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = FONT_NAME
    font.size = Pt(size)
    font.bold = bold
    if color:
        font.color.rgb = hex_to_rgb(color)
    if align is not None:
        paragraph.alignment = align


def remove_cell_borders(cell) -> None:
    try:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        borders = tcPr.find(qn("a:tcBorders"))
        if borders is None:
            borders = OxmlElement("a:tcBorders")
            tcPr.append(borders)
        for side in ("lnL", "lnR", "lnT", "lnB", "lnTlToBr", "lnBlToTr"):
            el = borders.find(qn(f"a:{side}"))
            if el is None:
                el = OxmlElement(f"a:{side}")
                borders.append(el)
            el.set("w", "0")
            el.clear()
            nofill = OxmlElement("a:noFill")
            el.append(nofill)
    except Exception:
        pass


def remove_table_gridlines(table) -> None:
    for r in table.rows:
        for c in r.cells:
            remove_cell_borders(c)


def add_background(slide, prs: Presentation, bg_source: Optional[Union[str, io.BytesIO]] = None):
    if bg_source is None:
        return
    try:
        slide.shapes.add_picture(bg_source, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    except Exception:
        pass  # Silently skip if background cannot be added


def add_cover_slide(prs: Presentation, title: str, bg_source: Optional[Union[str, io.BytesIO]] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Add background image covering the entire slide
    add_background(slide, prs, bg_source)

    left = Inches(0.5)
    top = Inches(2.4)
    width = prs.slide_width - Inches(1.0)
    height = Inches(2.0)
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    set_paragraph(p, title, 44, True, TITLE_COLOR, PP_ALIGN.CENTER)


def add_carteira_slide(prs: Presentation, nome: str, parsed: Dict, bg_source: Optional[Union[str, io.BytesIO]] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Add background image covering the entire slide
    add_background(slide, prs, bg_source)

    margin = Inches(0.5)
    gap = Inches(GAP_IN)
    content_w = prs.slide_width - 2 * margin

    # Columns per your layout image:
    # - Left column: Descrição, Valores mínimos, Distribuição por setor
    # - Right column: Composição da carteira (entire right side, top to bottom)
    left_w = Inches(LEFT_COL_W_IN)
    right_w = content_w - left_w - gap
    right_left = margin + left_w + gap

    # Title (full width), colored
    tb_title = slide.shapes.add_textbox(
        margin, Inches(0.2), content_w, Inches(0.7)
    )
    p = tb_title.text_frame.paragraphs[0]
    set_paragraph(p, nome, 28, True, TITLE_COLOR)

    # Left column: Descrição (colored subheader)
    desc_top = Inches(0.9)
    # Aumentei um pouco a altura para aliviar o conteúdo
    desc_h = Inches(1.55)  # antes 1.4
    tb_desc = slide.shapes.add_textbox(margin, desc_top, left_w, desc_h)
    tf = tb_desc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_paragraph(p, "Descrição", 15, True, TITLE_COLOR)
    p = tf.add_paragraph()
    set_paragraph(p, parsed.get("descricao") or "—", 11)

    # Left column: Valores mínimos (3 boxes) — descer um pouco
    # padding extra de 0.30" após a descrição (antes era 0.20")
    box_top = desc_top + desc_h + Inches(0.70)  # +0.10"
    box_gap = Inches(0.15)
    box_w = (left_w - 2 * box_gap) / 3
    box_h = Inches(0.7)
    labels = [
        ("Valor mínimo de aporte", parsed.get("aporte") or "—"),
        ("Valor mínimo de permanência", parsed.get("perman") or "—"),
        ("Valor mínimo de movimentação", parsed.get("mov") or "—"),
    ]
    for i, (lbl, val) in enumerate(labels):
        left = margin + i * (box_w + box_gap)
        tb = slide.shapes.add_textbox(left, box_top, box_w, box_h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        set_paragraph(p, lbl, 10, True, TITLE_COLOR)
        p = tf.add_paragraph()
        set_paragraph(p, val, 11)

    # Left column: Distribuição por setor — descer um pouco também
    # antes: set_top = box_top + box_h + Inches(0.30)
    set_top = box_top + box_h + Inches(0.80)  # +0.10"
    tb_hdr_set = slide.shapes.add_textbox(
        margin, set_top - Inches(0.32), left_w, Inches(0.3)
    )
    p = tb_hdr_set.text_frame.paragraphs[0]
    set_paragraph(p, "Distribuição por setor", 15, True, TITLE_COLOR)

    setores = parsed.get("setores", [])
    hdr_h = Inches(HDR_ROW_H_IN)
    row_h = Inches(BODY_ROW_H_IN)
    set_rows = 1 + max(1, len(setores))
    set_table_h = hdr_h + row_h * (set_rows - 1)

    set_shape = slide.shapes.add_table(
        set_rows, 2, margin, set_top, left_w, set_table_h
    )
    set_table = set_shape.table

    set_table.rows[0].height = hdr_h
    for i in range(1, set_rows):
        set_table.rows[i].height = row_h

    peso_w2 = Inches(0.9)
    setor_w = left_w - peso_w2
    set_table.columns[0].width = setor_w
    set_table.columns[1].width = peso_w2

    header_rgb = hex_to_rgb(TITLE_COLOR)

    for j, h in enumerate(["Setor", "Peso"]):
        cell = set_table.cell(0, j)
        cell.text = h
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.name = FONT_NAME
        run.font.size = Pt(HDR_FONT_SIZE)
        run.font.bold = True
        run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_rgb

    if setores:
        for i, (cat, pct) in enumerate(setores, start=1):
            vals = [cat, f"{pct}%"]
            is_odd = (i % 2 == 1)
            for j, val in enumerate(vals):
                cell = set_table.cell(i, j)
                cell.text = val
                tf = cell.text_frame
                tf.word_wrap = True
                run = tf.paragraphs[0].runs[0]
                run.font.name = FONT_NAME
                run.font.size = Pt(BODY_FONT_SIZE)
                cell.fill.solid()
                if is_odd:
                    cell.fill.fore_color.rgb = WHITE
                    cell.fill.transparency = 1.0
                else:
                    cell.fill.fore_color.rgb = LIGHT_GRAY
                    cell.fill.transparency = 0.0
    else:
        for j, val in enumerate(["—", "—"]):
            cell = set_table.cell(1, j)
            cell.text = val
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = FONT_NAME
            run.font.size = Pt(BODY_FONT_SIZE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            cell.fill.transparency = 1.0  # Transparent for the empty row (odd)

    remove_table_gridlines(set_table)

    # Right column: Composição da carteira (inteiro à direita, começa no topo)
    tbl_top = Inches(0.9)
    tb_hdr_tbl = slide.shapes.add_textbox(
        right_left, tbl_top - Inches(0.32), right_w, Inches(0.3)
    )
    p = tb_hdr_tbl.text_frame.paragraphs[0]
    set_paragraph(p, "Composição da carteira", 15, True, TITLE_COLOR)

    pos_data = parsed.get("carteira", [])
    pos_rows = 1 + max(1, len(pos_data))
    pos_table_h = hdr_h + row_h * (pos_rows - 1)

    pos_shape = slide.shapes.add_table(
        pos_rows, 3, right_left, tbl_top, right_w, pos_table_h
    )
    pos_table = pos_shape.table

    pos_table.rows[0].height = hdr_h
    for i in range(1, pos_rows):
        pos_table.rows[i].height = row_h

    peso_w = Inches(1.0)
    ticker_w = Inches(1.8)
    empresa_w = right_w - ticker_w - peso_w
    pos_table.columns[0].width = empresa_w
    pos_table.columns[1].width = ticker_w
    pos_table.columns[2].width = peso_w

    for j, h in enumerate(["Empresa", "Ticker", "Peso"]):
        cell = pos_table.cell(0, j)
        cell.text = h
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.name = FONT_NAME
        run.font.size = Pt(HDR_FONT_SIZE)
        run.font.bold = True
        run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_rgb

    for i in range(1, pos_rows):
        if i <= len(pos_data):
            empresa, ticker, peso = pos_data[i - 1]
            vals = [empresa, ticker, f"{peso}%"]
        else:
            vals = ["—", "—", "—"]
        is_odd = (i % 2 == 1)
        for j, val in enumerate(vals):
            cell = pos_table.cell(i, j)
            cell.text = val
            tf = cell.text_frame
            tf.word_wrap = True
            run = tf.paragraphs[0].runs[0]
            run.font.name = FONT_NAME
            run.font.size = Pt(BODY_FONT_SIZE)
            cell.fill.solid()
            if is_odd:
                cell.fill.fore_color.rgb = WHITE
                cell.fill.transparency = 1.0
            else:
                cell.fill.fore_color.rgb = LIGHT_GRAY
                cell.fill.transparency = 0.0

    remove_table_gridlines(pos_table)


def build_pptx(nome_fundo: str, carteiras: List[Dict], bg_source: Optional[Union[str, io.BytesIO]] = None) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # Fallback to local bg.png if no uploaded source provided
    if bg_source is None:
        script_dir = os.path.dirname(os.path.abspath(__file__))
        local_bg_path = os.path.join(script_dir, "assets", "bg.png")
        if os.path.exists(local_bg_path):
            bg_source = local_bg_path

    if nome_fundo.strip():
        add_cover_slide(prs, nome_fundo.strip(), bg_source)

    for c in carteiras:
        add_carteira_slide(prs, c["nome"], c["parsed"], bg_source)

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


# -------------------- Streamlit helpers --------------------


def style_table(df: pd.DataFrame):
    # Header color same as title; zebra stripes (Streamlit preview only)
    styler = df.style.set_table_styles(
        [
            {
                "selector": "th",
                "props": [
                    ("background-color", TITLE_COLOR),
                    ("color", "white"),
                    ("font-weight", "bold"),
                ],
            }
        ]
    )

    def zebra(row):
        color = "#F2F2F2" if row.name % 2 == 0 else "#FFFFFF"
        return [f"background-color: {color}"] * len(row)

    return styler.apply(zebra, axis=1)


# -------------------- Streamlit UI --------------------

st.set_page_config(page_title="Carteiras + PPTX", layout="wide")
st.title("Parser de Carteiras + Exportação PPTX")

st.markdown(
    "- Insira o Nome do Fundo (título principal do PPT).\n"
    "- Adicione carteiras: nome + texto no formato especificado.\n"
    "- Pré-visualize e gere o PowerPoint (16:9)."
)

if "carteiras" not in st.session_state:
    st.session_state.carteiras = []  # list of dicts: {nome, parsed}

nome_fundo = st.text_input("Nome do Fundo (capa do PPT)", value="")

# Background image uploader
uploaded_bg = st.file_uploader("Upload Background Image (PNG/JPG)", type=['png', 'jpg', 'jpeg'])
if uploaded_bg:
    st.info(f"Background image uploaded: {uploaded_bg.name}")
else:
    st.info("No background uploaded. Will use default 'assets/bg.png' if available, otherwise no background.")

st.divider()

st.subheader("Adicionar carteira")
c1, c2 = st.columns([1, 3])
with c1:
    nome_carteira = st.text_input("Nome da Carteira")
with c2:
    sample_hint = (
        "Cole o bloco começando em 'Descrição' e contendo as seções "
        "'Valor mínimo...', 'Composição da carteira' e 'Distribuição por setor'. "
        "Atenção: na distribuição, o percentual vem ANTES do setor."
    )
    texto = st.text_area("Texto da carteira", height=220, placeholder=sample_hint)

col_btn = st.columns([1, 1, 6])
with col_btn[0]:
    if st.button("Adicionar à lista", type="primary"):
        if not nome_carteira.strip():
            st.warning("Informe o Nome da Carteira.")
        elif not texto.strip():
            st.warning("Cole o texto da carteira.")
        else:
            try:
                parsed = parse_sections(texto)
                st.session_state.carteiras.append(
                    {"nome": nome_carteira.strip(), "parsed": parsed}
                )
                st.success(f"Carteira '{nome_carteira.strip()}' adicionada.")
                st.rerun()  # Use st.rerun() instead of experimental_rerun in newer Streamlit
            except Exception as e:
                st.error("Erro ao parsear a carteira.")
                st.exception(e)

with col_btn[1]:
    if st.button("Limpar lista de carteiras"):
        st.session_state.carteiras = []
        st.info("Lista limpa.")
        st.rerun()

st.divider()

st.subheader("Carteiras adicionadas")
if not st.session_state.carteiras:
    st.info("Nenhuma carteira adicionada ainda.")
else:
    for idx, c in enumerate(st.session_state.carteiras):
        parsed = c["parsed"]
        with st.expander(f"{idx + 1}. {c['nome']}", expanded=False):
            # Preview layout: left info + setores, right full positions
            left, right = st.columns([LEFT_COL_W_IN, 12 - LEFT_COL_W_IN], gap="small")

            with left:
                st.markdown("Descrição")
                st.write(parsed.get("descricao") or "—")

                c1, c2, c3 = st.columns(3)
                c1.metric("Aporte", parsed.get("aporte") or "—")
                c2.metric("Permanência", parsed.get("perman") or "—")
                c3.metric("Movimentação", parsed.get("mov") or "—")

                st.markdown("Distribuição por setor")
                setores = parsed.get("setores", [])
                if setores:
                    df_setores = pd.DataFrame(
                        setores, columns=["Setor", "Peso (%)"]
                    )
                    st.table(style_table(df_setores))
                else:
                    st.warning("Sem distribuição por setor.")

            with right:
                st.markdown("Composição da carteira")
                rows = parsed.get("carteira", [])
                if rows:
                    df_carteira = pd.DataFrame(
                        rows, columns=["Empresa", "Ticker", "Peso (%)"]
                    )
                    st.dataframe(
                        df_carteira, use_container_width=True, hide_index=True
                    )
                else:
                    st.warning("Tabela da carteira não identificada.")

            if st.button(f"Remover '{c['nome']}'", key=f"rm_{idx}"):
                st.session_state.carteiras.pop(idx)
                st.rerun()

st.divider()

if st.session_state.carteiras:
    ppt_bytes = build_pptx(nome_fundo, st.session_state.carteiras, uploaded_bg)
    st.download_button(
        label="Baixar PowerPoint (.pptx)",
        data=ppt_bytes,
        file_name="carteiras.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
else:
    st.info("Adicione pelo menos uma carteira para exportar o PPTX.")