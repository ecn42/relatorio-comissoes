import io
import re
from typing import Dict, List, Optional, Tuple

import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

# -------------------- Constants --------------------

TITLE_COLOR = "#8C6239"
FONT_NAME = "Montserrat"  # install for best PPT fidelity

LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# PPT sizing and styles
HDR_FONT_SIZE = 10
BODY_FONT_SIZE = 9
HDR_ROW_H_IN = 0.28
BODY_ROW_H_IN = 0.24

MARGIN_IN = 0.5


# -------------------- Utilities --------------------

def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.strip().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def normalize_text(text: str) -> str:
    return (
        text.replace("\u00A0", " ")
        .replace("\r\n", "\n")
        .replace("\r", "\n")
        .strip()
    )


def sanitize_line(line: str) -> str:
    return re.sub(r"\s+", " ", line.strip())


def is_header_line(line: str) -> bool:
    s = sanitize_line(line).lower()
    if not s:
        return False
    banned_exact = {
        "composição da carteira",
        "ativos",
        "peso retorno 12m %cdi 12m",
        "peso\tretorno 12m\t%cdi 12m",
    }
    if s in banned_exact:
        return False
    has_letters = bool(re.search(r"[A-Za-zÀ-ÿ]", s))
    has_percent = "%" in s
    return has_letters and not has_percent


def parse_pct_tokens_from_end(
    line: str,
) -> Optional[Tuple[str, str, str, str]]:
    """
    Return (name, peso%, ret12m%, cdi12m%) by taking the last three % tokens
    on the line; robust to extra % inside the name (e.g., '97,50% do CDI').
    """
    it = list(re.finditer(r"(\d[\d\.,]*)\s*%", line))
    if len(it) < 3:
        return None
    a, b, c = it[-3], it[-2], it[-1]
    name = line[: a.start()].strip()
    peso = line[a.start() : a.end()].strip()
    ret12 = line[b.start() : b.end()].strip()
    cdi = line[c.start() : c.end()].strip()

    def clean(p: str) -> str:
        p = p.replace(" ", "")
        if not p.endswith("%"):
            p = p + "%"
        return p

    return name, clean(peso), clean(ret12), clean(cdi)


def parse_top_composicao(
    lines: List[str], start_idx: int
) -> Tuple[int, List]:
    """
    Parse lines after 'Composição da Carteira' like:
      Pos Fixado    50,00%
    Stops on blank or first break after finding at least 1 row.
    Returns next index and list of (classe, peso_float, peso_str)
    """
    comp = []
    i = start_idx + 1
    while i < len(lines):
        raw = lines[i].strip()
        if not raw:
            if comp:
                break
            i += 1
            continue
        m = re.match(r"(.+?)\s+([\d\.,]+)\s*%$", raw)
        if not m:
            parts = re.split(r"\t+", raw)
            if len(parts) >= 2 and re.match(r"[\d\.,]+%$", parts[-1].strip()):
                classe = " ".join(parts[:-1]).strip(":- ")
                pct_str = parts[-1].strip()
            else:
                if comp:
                    break
                i += 1
                continue
        else:
            classe = m.group(1).strip(":- ")
            pct_str = m.group(2).strip() + "%"

        try:
            pct_float = float(
                pct_str.strip().rstrip("%").replace(".", "").replace(",", ".")
            )
        except Exception:
            pct_float = None

        comp.append((classe, pct_float, pct_str))
        i += 1
    return i, comp


def guess_portfolio_name(lines: List[str]) -> Optional[str]:
    for ln in lines:
        s = sanitize_line(ln)
        if not s:
            continue
        if s.lower() == "composição da carteira":
            break
        return s
    return None


def parse_portfolio_text(raw_text: str) -> Dict:
    """
    Parses text like the example provided.
    Returns:
      {
        'titulo_guess': Optional[str],
        'composicao': List[Tuple[str, Optional[float], str]],
        'holdings': List[Dict[str, str]]
      }
    """
    text = normalize_text(raw_text)
    lines = text.split("\n")

    idx_comp = -1
    for i, ln in enumerate(lines):
        if sanitize_line(ln).lower() == "composição da carteira":
            idx_comp = i
            break

    composicao = []
    if idx_comp != -1:
        _, composicao = parse_top_composicao(lines, idx_comp)

    holdings = []
    i = 0
    while i < len(lines):
        ln = lines[i].strip()
        if not ln:
            i += 1
            continue

        low = sanitize_line(ln).lower()
        if low in {
            "composição da carteira",
            "ativos",
            "peso retorno 12m %cdi 12m",
        }:
            i += 1
            continue

        if is_header_line(ln):
            section = sanitize_line(ln)
            j = i + 1
            while j < len(lines):
                s = sanitize_line(lines[j]).lower()
                if s in {"ativos", "peso retorno 12m %cdi 12m"}:
                    j += 1
                    continue
                break

            while j < len(lines):
                row = lines[j].strip()
                if not row:
                    j += 1
                    break
                if is_header_line(row) or (
                    sanitize_line(row).lower() == "composição da carteira"
                ):
                    break

                parsed = parse_pct_tokens_from_end(row)
                if parsed:
                    name, peso, ret12, cdi = parsed
                    holdings.append(
                        {
                            "grupo": section,
                            "ativo": sanitize_line(name),
                            "peso": peso,
                            "ret12m": ret12,
                            "cdi12m": cdi,
                        }
                    )
                j += 1
            i = j
            continue

        i += 1

    titulo_guess = guess_portfolio_name(lines)
    return {
        "titulo_guess": titulo_guess,
        "composicao": composicao,
        "holdings": holdings,
    }


# -------------------- PPTX helpers --------------------

def set_paragraph(
    paragraph,
    text: str,
    size: int,
    bold: bool = False,
    color: Optional[str] = None,
    align: Optional[int] = None,
):
    if hasattr(paragraph, "clear"):
        paragraph.clear()
    else:
        paragraph.text = ""
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = FONT_NAME
    font.size = Pt(size)
    font.bold = bold
    if color:
        font.color.rgb = hex_to_rgb(color)
    if align is not None:
        paragraph.alignment = align


def remove_cell_borders(cell) -> None:
    try:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        borders = tcPr.find(qn("a:tcBorders"))
        if borders is None:
            borders = OxmlElement("a:tcBorders")
            tcPr.append(borders)
        for side in ("lnL", "lnR", "lnT", "lnB", "lnTlToBr", "lnBlToTr"):
            el = borders.find(qn(f"a:{side}"))
            if el is None:
                el = OxmlElement(f"a:{side}")
                borders.append(el)
            el.set("w", "0")
            el.clear()
            nofill = OxmlElement("a:noFill")
            el.append(nofill)
    except Exception:
        pass


def remove_table_gridlines(table) -> None:
    for r in table.rows:
        for c in r.cells:
            remove_cell_borders(c)


def add_cover_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    left = Inches(MARGIN_IN)
    top = Inches(3.2)  # centered better on portrait
    width = prs.slide_width - Inches(2 * MARGIN_IN)
    height = Inches(1.8)
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    set_paragraph(p, title, 44, True, TITLE_COLOR, PP_ALIGN.CENTER)


def add_portfolio_slide(prs: Presentation, nome: str, parsed: Dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    margin = Inches(MARGIN_IN)
    content_w = prs.slide_width - 2 * margin

    # Title
    title_top = Inches(0.2)
    tb_title = slide.shapes.add_textbox(
        margin, title_top, content_w, Inches(0.7)
    )
    p = tb_title.text_frame.paragraphs[0]
    set_paragraph(p, nome, 28, True, TITLE_COLOR)

    header_rgb = hex_to_rgb(TITLE_COLOR)
    hdr_h = Inches(HDR_ROW_H_IN)
    row_h = Inches(BODY_ROW_H_IN)

    # Composição (full width, on top)
    comp_table_top = title_top + Inches(0.7) + Inches(0.4)
    tb_hdr_comp = slide.shapes.add_textbox(
        margin, comp_table_top - Inches(0.32), content_w, Inches(0.3)
    )
    p = tb_hdr_comp.text_frame.paragraphs[0]
    set_paragraph(p, "Composição da Carteira", 15, True, TITLE_COLOR)

    composicao = parsed.get("composicao", [])
    comp_rows = 1 + max(1, len(composicao))
    comp_table_h = hdr_h + row_h * (comp_rows - 1)

    comp_shape = slide.shapes.add_table(
        comp_rows, 2, margin, comp_table_top, content_w, comp_table_h
    )
    comp_table = comp_shape.table

    comp_table.rows[0].height = hdr_h
    for i in range(1, comp_rows):
        comp_table.rows[i].height = row_h

    peso_w = Inches(1.0)
    classe_w = content_w - peso_w
    comp_table.columns[0].width = classe_w
    comp_table.columns[1].width = peso_w

    # Header
    for j, h in enumerate(["Classe", "Peso"]):
        cell = comp_table.cell(0, j)
        cell.text = h
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.name = FONT_NAME
        run.font.size = Pt(HDR_FONT_SIZE)
        run.font.bold = True
        run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_rgb

    if composicao:
        for i, (classe, _pct_float, pct_str) in enumerate(composicao, start=1):
            vals = [classe, pct_str]
            is_even = (i % 2 == 0)
            bg = LIGHT_GRAY if is_even else WHITE
            for j, val in enumerate(vals):
                cell = comp_table.cell(i, j)
                cell.text = val
                tf = cell.text_frame
                tf.word_wrap = True
                par = tf.paragraphs[0]
                run = par.runs[0]
                run.font.name = FONT_NAME
                run.font.size = Pt(BODY_FONT_SIZE)
                if j == 1:
                    par.alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
    else:
        for j, val in enumerate(["—", "—"]):
            cell = comp_table.cell(1, j)
            cell.text = val
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = FONT_NAME
            run.font.size = Pt(BODY_FONT_SIZE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

    remove_table_gridlines(comp_table)

    # Ativos (full width, below)
    assets_top = comp_table_top + comp_table_h + Inches(0.6)
    tb_hdr_tbl = slide.shapes.add_textbox(
        margin, assets_top - Inches(0.32), content_w, Inches(0.3)
    )
    p = tb_hdr_tbl.text_frame.paragraphs[0]
    set_paragraph(p, "Ativos da carteira", 15, True, TITLE_COLOR)

    rows = parsed.get("holdings", [])
    pos_rows = 1 + max(1, len(rows))
    pos_table_h = hdr_h + row_h * (pos_rows - 1)

    pos_shape = slide.shapes.add_table(
        pos_rows, 5, margin, assets_top, content_w, pos_table_h
    )
    pos_table = pos_shape.table

    pos_table.rows[0].height = hdr_h
    for i in range(1, pos_rows):
        pos_table.rows[i].height = row_h

    # Column widths on portrait width
    grupo_w = Inches(1.6)
    peso_w = Inches(0.9)
    ret_w = Inches(1.0)
    cdi_w = Inches(0.9)
    ativo_w = content_w - (grupo_w + peso_w + ret_w + cdi_w)

    pos_table.columns[0].width = grupo_w
    pos_table.columns[1].width = ativo_w
    pos_table.columns[2].width = peso_w
    pos_table.columns[3].width = ret_w
    pos_table.columns[4].width = cdi_w

    headers = ["Grupo", "Ativo", "Peso", "Ret. 12M", "%CDI 12M"]
    for j, h in enumerate(headers):
        cell = pos_table.cell(0, j)
        cell.text = h
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.name = FONT_NAME
        run.font.size = Pt(HDR_FONT_SIZE)
        run.font.bold = True
        run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_rgb

    for i in range(1, pos_rows):
        if i <= len(rows):
            r = rows[i - 1]
            vals = [r["grupo"], r["ativo"], r["peso"], r["ret12m"], r["cdi12m"]]
        else:
            vals = ["", "", "", "", ""]
        is_even = (i % 2 == 0)
        bg = LIGHT_GRAY if is_even else WHITE
        for j, val in enumerate(vals):
            cell = pos_table.cell(i, j)
            cell.text = val
            tf = cell.text_frame
            tf.word_wrap = True
            par = tf.paragraphs[0]
            run = par.runs[0]
            run.font.name = FONT_NAME
            run.font.size = Pt(BODY_FONT_SIZE)
            if j in (2, 3, 4):
                par.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    remove_table_gridlines(pos_table)


def build_pptx(nome_doc: str, carteiras: List[Dict]) -> io.BytesIO:
    prs = Presentation()

    # Portrait "widescreen": rotate 16:9 to vertical
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(13.333)

    if nome_doc.strip():
        add_cover_slide(prs, nome_doc.strip())

    for c in carteiras:
        add_portfolio_slide(prs, c["nome"], c["parsed"])

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


# -------------------- Streamlit helpers --------------------

def style_table(df: pd.DataFrame):
    styler = df.style.set_table_styles(
        [
            {
                "selector": "th",
                "props": [
                    ("background-color", TITLE_COLOR),
                    ("color", "white"),
                    ("font-weight", "bold"),
                ],
            }
        ]
    )

    def zebra(row):
        color = "#F2F2F2" if row.name % 2 == 0 else "#FFFFFF"
        return [f"background-color: {color}"] * len(row)

    return styler.apply(zebra, axis=1)


# -------------------- Streamlit UI --------------------

st.set_page_config(page_title="Portfólios (Vertical) + PPTX", layout="centered")
st.title("Parser de Portfólios (Texto) + PPTX Vertical 16:9")

st.markdown(
    "- Cole o texto no formato do exemplo (Composição da Carteira no topo, "
    "seções com Ativos, Peso, Retorno 12M, %CDI 12M).\n"
    "- Adicione vários portfólios e gere um PowerPoint em retrato (vertical), "
    "com composição em cima e ativos abaixo.\n"
    "- Cores e estilo iguais ao exemplo."
)

if "carteiras" not in st.session_state:
    st.session_state.carteiras = []  # list of dicts: {nome, parsed}

nome_doc = st.text_input("Nome do Documento (capa do PPT)", value="")

st.divider()

st.subheader("Adicionar portfólio")
c1, c2 = st.columns([1, 3])
with c1:
    nome_port = st.text_input("Nome do Portfólio (opcional)")
with c2:
    placeholder = (
        "Exemplo:\n"
        "ASSESSORIA DIGITAL - CONSERVADOR\n"
        "Composição da Carteira\n"
        "Pos Fixado\t50,00%\n"
        "Inflação\t20,00%\n"
        "Multimercados\t17,50%\n"
        "Previdencia\t12,50%\n\n"
        "Renda Fixa POS FIXADOS\n"
        "Peso\tRetorno 12M\t%CDI 12M\n"
        "Ativos\n"
        "CDB - Banco Luso - 104.0 %CDI\t7,50%\t14,27%\t104,27%\n"
        "CDB - Neon Financeira - 106.75 %CDI\t7,50%\t14,68%\t107,22%\n"
        "CRA - GRUPO JOSÉ ALVES - CDI + 0,85%\t2,50%\t14,65%\t107,30%\n"
        "CRA - MINERVA - 97,50% do CDI\t2,50%\t13,32%\t97,34%\n\n"
        "Renda Fixa INFLAÇÃO Pos Fixados\n"
        "Peso\tRetorno 12M\t%CDI 12M\n"
        "Ativos\n"
        "RUMOB7 - RUMO S.A - IPCA + 6,70%\t6,00%\t11,95%\t87,30%\n"
        "ERDVB4 - ...\t7,00%\t12,23%\t89,36%\n"
    )
    texto = st.text_area("Texto do Portfólio", height=260, placeholder=placeholder)

col_btn = st.columns([1, 1, 6])
with col_btn[0]:
    if st.button("Adicionar à lista", type="primary"):
        if not texto.strip():
            st.warning("Cole o texto do portfólio.")
        else:
            try:
                parsed = parse_portfolio_text(texto)
                nome_final = (
                    nome_port.strip()
                    or parsed.get("titulo_guess")
                    or "Portfólio"
                )
                st.session_state.carteiras.append(
                    {"nome": nome_final, "parsed": parsed}
                )
                st.success(f"Portfólio '{nome_final}' adicionado.")
            except Exception as e:
                st.error("Erro ao parsear o portfólio.")
                st.exception(e)

with col_btn[1]:
    if st.button("Limpar lista de portfólios"):
        st.session_state.carteiras = []
        st.info("Lista limpa.")

st.divider()

st.subheader("Portfólios adicionados")
if not st.session_state.carteiras:
    st.info("Nenhum portfólio adicionado ainda.")
else:
    for idx, c in enumerate(st.session_state.carteiras):
        parsed = c["parsed"]
        with st.expander(f"{idx + 1}. {c['nome']}", expanded=False):
            st.markdown("Composição da Carteira")
            comp = parsed.get("composicao", [])
            if comp:
                df_comp = pd.DataFrame(
                    [(a, c_str) for a, _c, c_str in comp],
                    columns=["Classe", "Peso (%)"],
                )
                st.table(style_table(df_comp))
            else:
                st.warning("Sem composição no topo.")

            st.markdown("Ativos da carteira")
            rows = parsed.get("holdings", [])
            if rows:
                df_h = pd.DataFrame(
                    rows,
                    columns=["grupo", "ativo", "peso", "ret12m", "cdi12m"],
                ).rename(
                    columns={
                        "grupo": "Grupo",
                        "ativo": "Ativo",
                        "peso": "Peso (%)",
                        "ret12m": "Retorno 12M",
                        "cdi12m": "%CDI 12M",
                    }
                )
                st.dataframe(df_h, use_container_width=True, hide_index=True)
            else:
                st.warning("Nenhum ativo identificado.")

            if st.button(f"Remover '{c['nome']}'", key=f"rm_{idx}"):
                st.session_state.carteiras.pop(idx)
                st.experimental_rerun()

st.divider()

if st.session_state.carteiras:
    ppt_bytes = build_pptx(nome_doc, st.session_state.carteiras)
    st.download_button(
        label="Baixar PowerPoint (.pptx)",
        data=ppt_bytes,
        file_name="portfolios_vertical.pptx",
        mime=(
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        ),
    )
else:
    st.info("Adicione pelo menos um portfólio para exportar o PPTX.")